import colorama
colorama.deinit()
import threading

from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import os
import uuid
import traceback
import json
from datetime import datetime
import datetime as _dt
from flask_socketio import SocketIO, emit
import sys

class Logger:
    def __init__(self, original_stream, log_path):
        self.original_stream = original_stream
        self.log_path = log_path

    def write(self, data):
        try:
            # Safe terminal output: replace non-ASCII with ? for Windows compatibility
            self.original_stream.write(data.encode(sys.stdout.encoding, errors='replace').decode(sys.stdout.encoding))
            self.original_stream.flush()
        except Exception:
            try:
                self.original_stream.write(data.encode('ascii', errors='replace').decode('ascii'))
            except Exception:
                pass

        try:
            with open(self.log_path, "a", encoding="utf-8") as f:
                f.write(data)
        except Exception:
            pass

    def flush(self):
        self.original_stream.flush()

sys.stdout = Logger(sys.stdout, "backend_debug.log")
sys.stderr = Logger(sys.stderr, "backend_debug.log")

print(f"\n--- BACKEND STARTED AT {datetime.now()} ---")


def get_next_monday() -> str:
    """
    Returns the date of the next (upcoming) Monday in MMDDYYYY format.
    If today is Monday, returns today.
    """
    today = _dt.date.today()
    days_until_monday = (7 - today.weekday()) % 7  # weekday(): Monday=0
    if days_until_monday == 0:
        target = today
    else:
        target = today + _dt.timedelta(days=days_until_monday)
    return target.strftime("%m%d%Y")

from utils.ppt_slide_summary import extract_slides_text, summarize_slides
from excel_processing.read_excel import read_visible_excel_sheets_to_json
from fix_ppt_fonts import fix_ppt_fonts
from ppt_generation.slides import create_ppt_from_excel
from ppt_generation import ppt_helper
from fix_ppt import fix_ppt
from sockets.events import register_socket_events

from dotenv import load_dotenv

from utils.process_slides_and_stream import process_slides_and_stream

load_dotenv()
import flask.cli
flask.cli.show_server_banner = lambda *x: None

app = Flask(__name__)
CORS(app)

# Use eventlet in production (gunicorn), threading in dev
async_mode = os.environ.get("ASYNC_MODE", "threading")
socketio = SocketIO(app, cors_allowed_origins="*", async_mode=async_mode)
register_socket_events(socketio)

UPLOAD_FOLDER = "./uploads"
GENERATED_FOLDER = "./generated"
EXTRACTED_FOLDER = "./extracted_json"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(GENERATED_FOLDER, exist_ok=True)
os.makedirs(EXTRACTED_FOLDER, exist_ok=True)

@app.route("/api/upload-report", methods=["POST"])
def upload_report():
    print("\n" + "="*50)
    print(">>> INCOMING REQUEST: /api/upload-report")
    print("="*50)
    # Support multiple excel files: frontend sends 'excel_files' (multiple)
    excel_files = request.files.getlist('excel_files')

    # Backward compatibility: also accept single 'excel' field
    if not excel_files:
        single_excel = request.files.get('excel')
        if single_excel and single_excel.filename:
            excel_files = [single_excel]

    # Filter out empty filenames
    excel_files = [f for f in excel_files if f.filename and f.filename != '']
    if not excel_files:
        return jsonify({"message": "At least one Excel file is required"}), 400

    unique_id = str(uuid.uuid4())
    # Use first excel name for download naming
    excel_name = os.path.splitext(excel_files[0].filename)[0]

    # Use hardcoded template from backend; fall back to uploaded PPT if provided
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))
    default_template = os.path.join(BASE_DIR, "templates", "default_template.pptx")

    ppt_file = request.files.get('ppt')
    if ppt_file and ppt_file.filename and ppt_file.filename != '':
        # User uploaded a custom template — use it
        ppt_upload_path = os.path.join(UPLOAD_FOLDER, f"{unique_id}.pptx")
        ppt_file.save(ppt_upload_path)
        print(f"Using custom template: {ppt_file.filename}")
    elif os.path.exists(default_template):
        ppt_upload_path = default_template
        print(f"Using default template: {default_template}")
    else:
        return jsonify({"message": "No PPT template found. Place default_template.pptx in the templates/ folder."}), 400

    # Capture date format from frontend (default UK)
    date_format = request.form.get("date_format", "UK")
    ppt_helper.DATE_FORMAT_REGION = date_format.upper()

    roles_json = request.form.get("roles")
    try:
        project_roles = json.loads(roles_json) if roles_json else {}
    except Exception:
        project_roles = {}
        print(" Invalid roles JSON received")

    print("\nReceived dynamic roles from frontend:")
    for k, v in project_roles.items():
        print(f"  {k} -> {v}")

    # Save all excel files with indexed names
    excel_upload_paths = []
    for idx, ef in enumerate(excel_files):
        excel_path = os.path.join(UPLOAD_FOLDER, f"{unique_id}_excel_{idx}.xlsx")
        ef.save(excel_path)
        excel_upload_paths.append({"path": excel_path, "name": ef.filename})
        print(f"Saved Excel [{idx}]: {ef.filename} -> {excel_path}")

    try:
        print("PPT path:", ppt_upload_path)
        print(f"Total Excel files: {len(excel_upload_paths)}")

        project_title = request.form.get("project_title", "")
        if not project_title or project_title.strip() == "":
            excel_filename_base = os.path.splitext(excel_files[0].filename)[0]
            clean_title = excel_filename_base

            # Remove known suffix terms that come after a "_"
            # e.g. "SB - Master Data Feed Redesign_Project Tracker" -> "SB - Master Data Feed Redesign"
            strip_terms = ["Project Tracker", "Financial Summary", "project tracker", "financial summary"]
            for term in strip_terms:
                suffix = f"_{term}"
                if suffix in clean_title:
                    clean_title = clean_title[:clean_title.rfind(suffix)]
                # Also remove as plain text (without underscore)
                clean_title = clean_title.replace(term, "")

            project_title = clean_title.strip().strip("_").strip()
            print(f"Extracted clean project title from filename: {project_title}")

        author = request.form.get("author", "")
        report_date = request.form.get("report_date", "")
        ppt_generated_path = os.path.join(GENERATED_FOLDER, f"{unique_id}_report.pptx")

        # Debug print confirmation
        print("\nFinal project_roles to send to PPT generation:")
        for role, name in project_roles.items():
            print(f"  {role} -> {name}")

        all_sheets_data = {}
        unified_excel_data = {}
        all_extracted_json_files = []

        for excel_info in excel_upload_paths:
            excel_upload_path = excel_info["path"]
            excel_filename = excel_info["name"]
            print(f"\n{'='*50}")
            print(f"Processing Excel: {excel_filename}")
            print(f"{'='*50}")

            # Step 1: Extract sheets as styled JSON files
            extracted_json_files = read_visible_excel_sheets_to_json(excel_upload_path, EXTRACTED_FOLDER)
            all_extracted_json_files.extend(extracted_json_files)
            
            for json_path in extracted_json_files:
                with open(json_path, "r", encoding="utf-8") as f:
                    sheet_json = json.load(f)
                
                filename_base = os.path.splitext(os.path.basename(json_path))[0]
                sheet_name = filename_base.split("==", 1)[1] if "==" in filename_base else filename_base
                
                # Global aggregate for summary slides
                all_sheets_data[filename_base] = sheet_json
                
                # Unified data for all slides
                unique_name = sheet_name
                counter = 2
                while unique_name in unified_excel_data:
                    unique_name = f"{sheet_name} ({counter})"
                    counter += 1
                unified_excel_data[unique_name] = sheet_json

        # Step 2: Generate PPT in one batch
        if not unified_excel_data:
            print("No data extracted from any Excel file.")
            return jsonify({"message": "No data found in uploaded Excel files"}), 400

        create_ppt_from_excel(
            ppt_upload_path,
            unified_excel_data,
            ppt_generated_path,
            project_title=project_title,
            project_roles=project_roles,
            author=author,
            skip_header_slides=False,
            summary_data=all_sheets_data
        )


        # Fix the final generated PPT
        output_path = fix_ppt(ppt_generated_path)

        # Only run AI summarization if user opted in (saves tokens)
        auto_summarize = request.form.get("auto_summarize", "false").lower() == "true"

        if auto_summarize:
            BASE_DIR = os.path.dirname(os.path.abspath(__file__))

            # Load summarization prompt
            prompt_file = os.path.join(BASE_DIR, "prompts", "summarization_prompt.txt")
            print("Looking for prompt at:", prompt_file)
            with open(prompt_file, "r", encoding="utf-8") as f:
                prompt_template = f.read()

            threading.Thread(
                target=process_slides_and_stream,
                args=(ppt_generated_path, prompt_template, unique_id, socketio)
            ).start()
        else:
            print("Auto-summarize disabled — skipping AI summary generation")

        if not os.path.exists(ppt_generated_path):
            return jsonify({"message": "Failed to create PPT file"}), 500

        # Build download name: {excel_base}_Status Report_{next_monday_MMDDYYYY}.pptx
        monday_date = get_next_monday()
        
        # Clean excel_name for the final filename
        clean_excel_name = excel_name
        strip_terms = ["Project Tracker", "Financial Summary", "project tracker", "financial summary"]
        for term in strip_terms:
            suffix = f"_{term}"
            if suffix in clean_excel_name:
                clean_excel_name = clean_excel_name[:clean_excel_name.rfind(suffix)]
            clean_excel_name = clean_excel_name.replace(term, "")
        clean_excel_name = clean_excel_name.strip().strip("_").strip()
        
        status_report_name = f"{clean_excel_name}_Status Report_{monday_date}"

        response = {
            "pptDownloadUrl": f"http://localhost:5000/api/download-report/{os.path.basename(output_path)}/ppt?name={status_report_name}",
            "extractedJson": [os.path.basename(j) for j in all_extracted_json_files],
            "new_fixed_ppt": os.path.basename(output_path),
            "friendly_name": f"{status_report_name}.pptx",
            "excel_count": len(excel_upload_paths),
        }
        return jsonify(response)

    except Exception as e:
        print("\n" + "!"*50)
        print("CRITICAL ERROR IN UPLOAD_REPORT:")
        traceback.print_exc()
        print("!"*50 + "\n")
        sys.stdout.flush()
        return jsonify({
            "message": "Error processing report",
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500


@app.route("/api/download-report/<file_id>/<file_type>", methods=["GET"])
def download_report(file_id, file_type):
    if file_type == "ppt":
        file_path = os.path.join(GENERATED_FOLDER, f"{file_id}_report.pptx")
        
        # Get excel name from query param (default = Report)
        excel_name = request.args.get("name", "Report")
        # Final download name: <ExcelFileName>.pptx
        download_name = f"{excel_name}.pptx"
        
        # download_name = "PptReport.pptx"
        mimetype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    else:
        return jsonify({"message": "Invalid file type"}), 400

    if not os.path.exists(file_path):
        return jsonify({"message": "Report not found"}), 404

    return send_file(file_path, as_attachment=True, download_name=download_name, mimetype=mimetype)


@app.route("/api/download-report/<file_name>", methods=["GET"])
def download(file_name):
    ppt_path = os.path.join("generated", file_name)
    if not os.path.exists(ppt_path):
        print("PPT file not found:", ppt_path)
        return jsonify({"message": "PPT file not found"}), 404
    print("Sending file:", ppt_path)
    return send_file(ppt_path, as_attachment=True)


@app.route("/api/update-project-details", methods=["POST"])
def update_project_details():
    """
    Update project details (title, author, date, roles, distribution) in an
    already-generated PPT file. This is called from the Sidebar 'Details' tab
    so users can fix mistakes without re-uploading everything.
    """
    from pptx import Presentation
    from pptx.util import Pt
    from datetime import datetime as dt

    try:
        data = request.get_json()
        ppt_name = data.get("ppt_name")
        if not ppt_name:
            return jsonify({"message": "ppt_name is required"}), 400

        ppt_path = os.path.join("generated", ppt_name)
        if not os.path.exists(ppt_path):
            return jsonify({"message": f"PPT file not found: {ppt_name}"}), 404

        project_title = data.get("project_title", "")
        author = data.get("author", "")
        report_date = data.get("report_date", "")
        date_format = data.get("date_format", "UK").upper()
        roles = data.get("roles", {})

        prs = Presentation(ppt_path)

        def _set_para_text(para, text):
            """Set paragraph text, clearing all runs and putting text in run 0."""
            for run in para.runs:
                run.text = ""
            if para.runs:
                para.runs[0].text = text
            else:
                para.text = text

        # ── Update Slide 1 ──
        slide1 = prs.slides[0]

        for shape in slide1.shapes:
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip()

                # Update project title + date + author block (Rectangle 6 typically)
                if "Weekly Project Status Report" in full_text or "Status Report" in full_text:
                    paras = shape.text_frame.paragraphs

                    # Para 0: Project Title
                    if project_title and len(paras) > 0:
                        _set_para_text(paras[0], project_title)

                    # Para 1: "Weekly Project Status Report" — leave as-is

                    # Para 2: Date line
                    if report_date and len(paras) > 2:
                        try:
                            d = dt.strptime(report_date, "%Y-%m-%d")
                            if date_format == "US":
                                date_str = d.strftime("%B %d, %Y")
                            else:
                                day = d.day
                                suffix = "th" if 11 <= day <= 13 else {1: "st", 2: "nd", 3: "rd"}.get(day % 10, "th")
                                date_str = f"{day}{suffix} {d.strftime('%B %Y')}"
                        except Exception:
                            date_str = report_date
                        _set_para_text(paras[2], f"Date:  {date_str}")

                    # Para 3: Author line
                    if author and len(paras) > 3:
                        # Author para typically has 2 runs: "Author: " + "name"
                        if len(paras[3].runs) >= 2:
                            paras[3].runs[0].text = "Author: "
                            paras[3].runs[1].text = author
                            # Clear any extra runs
                            for r in paras[3].runs[2:]:
                                r.text = ""
                        else:
                            _set_para_text(paras[3], f"Author: {author}")

            # Update Distribution table
            if shape.has_table:
                tbl = shape.table
                first_cell = tbl.cell(0, 0).text.strip().lower()
                if "distribution" in first_cell:
                    print(f"Updating distribution table: {len(tbl.rows)} rows")
                    for r in range(1, len(tbl.rows)):
                        role_label = tbl.cell(r, 1).text.strip()
                        if role_label in roles:
                            person_name = roles[role_label]
                            # Set name in column 0
                            cell = tbl.cell(r, 0)
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    run.text = ""
                                if para.runs:
                                    para.runs[0].text = person_name
                                else:
                                    para.text = person_name
                            print(f"  {role_label} -> {person_name}")

        # ── Update Slide 2 (roles table, project start date) ──
        if len(prs.slides) > 1:
            slide2 = prs.slides[1]

            # Update project title in slide 2 header if present
            if project_title:
                for shape in slide2.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            txt = para.text.strip()
                            # Look for the status title text
                            if "project status" in txt.lower() or "status report" in txt.lower():
                                for run in para.runs:
                                    run.text = ""
                                if para.runs:
                                    para.runs[0].text = f"{project_title} - Project Status"

            # Update roles table (Table 15 - "SWB Head of Customer Solutions")
            for shape in slide2.shapes:
                if shape.has_table:
                    tbl = shape.table
                    for r in range(len(tbl.rows)):
                        for c in range(len(tbl.columns)):
                            cell_text = tbl.cell(r, c).text.strip()
                            # Match role labels and update adjacent cells
                            if cell_text.endswith(":") and cell_text[:-1] in roles:
                                # The value is in the next column
                                if c + 1 < len(tbl.columns):
                                    val_cell = tbl.cell(r, c + 1)
                                    person = roles[cell_text[:-1]]
                                    for para in val_cell.text_frame.paragraphs:
                                        for run in para.runs:
                                            run.text = ""
                                        if para.runs:
                                            para.runs[0].text = person
                                        else:
                                            para.text = person

        prs.save(ppt_path)
        print(f"\nProject details updated in {ppt_name}")
        print(f"  Title: {project_title}")
        print(f"  Author: {author}")
        print(f"  Date: {report_date} ({date_format})")
        print(f"  Roles: {len(roles)} entries")

        return jsonify({"message": "Project details updated successfully"})

    except Exception as e:
        print("\n" + "!"*50)
        print("ERROR IN UPDATE_PROJECT_DETAILS:")
        traceback.print_exc()
        print("!"*50 + "\n")
        sys.stdout.flush()
        return jsonify({
            "message": "Error updating project details",
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8000))
    debug = os.environ.get("FLASK_ENV", "development") == "development"
    socketio.run(
        app,
        host="0.0.0.0",
        port=port,
        debug=debug,
        use_reloader=debug,
        allow_unsafe_werkzeug=True
    )