from pptx import Presentation

from ai_helper.open_ai_helper import OpenAIHelper


def extract_slides_text(ppt_path):
    prs = Presentation(ppt_path)
    slides_text = []

    for idx, slide in enumerate(prs.slides):
        text_content = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

        full_text = "\n".join(text_content).strip()

        slides_text.append({
            "slide_number": idx + 1,
            "text": full_text
        })

    return slides_text


def summarize_slides(slides_text, prompt_template):
    """
    Summarizes each slide's text using OpenAI.
    Creates its own OpenAIHelper instance — no external 'ai' dependency needed.
    """
    ai = OpenAIHelper()
    summaries = []

    for slide in slides_text:
        if not slide["text"].strip():
            continue

        prompt = prompt_template.replace("{{text}}", slide["text"])

        summary = ai.ask(
            prompt=prompt,
            system_prompt="You are a professional Project Manager. Summarize concisely.",
            max_tokens=80,
            temperature=0.3
        )

        summaries.append({
            "slide_number": slide["slide_number"],
            "summary": summary
        })

    return summaries
