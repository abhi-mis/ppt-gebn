"""
fix_ppt_fonts.py — Standalone post-generation PPT font & style fixer.

Usage:
    python fix_ppt_fonts.py report.pptx
    python fix_ppt_fonts.py report.pptx --config my_style.json

Output:
    Saves fixed file as <filename>-styled.pptx in the same directory.

What it fixes:
    - Table header rows: enforces bold, font name, font size, font color, background color
    - Table body rows: enforces font name, font size, font color
    - Section titles (textbox titles like "Key Milestones", "Risks"): font name, size, color, bold
    - Consistent alignment across all cells
    - Optionally strips None/inherit font values and sets explicit defaults

Configuration:
    Edit STYLE_CONFIG below or pass a JSON file via --config flag.
"""

import sys
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ═══════════════════════════════════════════════════════════════
# STYLE CONFIG — Edit these values to customize your PPT styling
# ═══════════════════════════════════════════════════════════════
STYLE_CONFIG = {
    # ── Table Header Row (first row of each table) ──
    "table_header": {
        "font_name": "Calibri",
        "font_size_pt": 14,
        "font_bold": True,
        "font_italic": False,
        "font_color": "000000",          # black text
        "bg_color": "BDD7EE",            # light blue background
        "alignment": "center",           # center | left | right
        "vertical_alignment": "middle",  # top | middle | bottom
    },

    # ── Table Body Rows (all rows after header) ──
    "table_body": {
        "font_name": "Calibri",
        "font_size_pt": 12,
        "font_bold": False,
        "font_italic": False,
        "font_color": "000000",          # black text
        "alignment": "left",             # center | left | right
        "vertical_alignment": "middle",
    },

    # ── Section Titles ("Key Milestones", "Risks", etc.) ──
    "section_title": {
        "font_name": "Tahoma",
        "font_size_pt": 20,
        "font_bold": True,
        "font_italic": False,
        "font_color": "14366B",          # dark blue
    },

    # ── RAG Legend (Red/Amber/Green text box) ──
    # Set to null/None to leave untouched
    "rag_legend": {
        "font_name": "Tahoma",
        "font_size_pt": 8.5,
    },

    # ── Slides to skip (1-indexed). Title, status, thank-you, appendix ──
    "skip_slides": [1, 2],

    # ── Also skip slides that contain only a single picture (thank-you slides) ──
    "skip_picture_only_slides": True,

    # ── Also skip slides with placeholder shapes (appendix slides) ──
    "skip_placeholder_slides": True,
}


# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════

def hex_to_rgb(hex_str):
    """Convert '14366B' or '#14366B' to RGBColor."""
    h = hex_str.replace("#", "").strip()
    if len(h) == 8:  # ARGB
        h = h[2:]
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def get_alignment(align_str):
    """Convert string alignment to PP_ALIGN enum."""
    m = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    return m.get(align_str.lower(), PP_ALIGN.LEFT)


def get_vertical_alignment(align_str):
    """Convert string to MSO_VERTICAL_ANCHOR enum."""
    m = {
        "top": MSO_VERTICAL_ANCHOR.TOP,
        "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
        "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
    }
    return m.get(align_str.lower(), MSO_VERTICAL_ANCHOR.MIDDLE)


# ── Shape identification (same logic as fix_ppt.py) ──

RAG_SHAPE_NAMES = {
    "Flowchart: Connector 13",
    "Flowchart: Connector 3",
    "Flowchart: Connector 5",
}
RAG_TEXTBOX_SNIPPET = "Project at risk"


def is_rag_legend(shape):
    if shape.name in RAG_SHAPE_NAMES:
        return True
    if shape.has_text_frame and RAG_TEXTBOX_SNIPPET in shape.text_frame.text:
        return True
    return False


def is_section_title(shape):
    """
    Section title text boxes: "Key Milestones", "Risks", "Deliverables (Contd..)", etc.
    """
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    if not text:
        return False
    if is_rag_legend(shape):
        return False
    # Exclude known non-title shapes
    if shape.name in ("TextBox 3", "Rectangle 6"):
        return False
    if "Project at risk" in text or "recover/mitigate" in text:
        return False
    if "Project on track" in text:
        return False
    if text.startswith("*See Milestones"):
        return False
    if "©" in text:
        return False
    # None-text boxes
    if text in ("• None", "•None", "None", "- None"):
        return False
    return True


def is_header_row(table, row_idx):
    """
    Detect if a row is a header row.
    Heuristics:
      - Row 0 is almost always a header
      - A row where ALL cells are bold is likely a header
      - A row with blue background fill is likely a header
    """
    if row_idx == 0:
        return True

    # Check if all cells in this row have bold text
    all_bold = True
    has_text = False
    for ci in range(len(table.columns)):
        cell = table.cell(row_idx, ci)
        text = cell.text.strip()
        if text:
            has_text = True
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    if run.font.bold is not True:
                        all_bold = False
                    break
                break

    if has_text and all_bold:
        return True

    # Check background color — BDD7EE (light blue) is header
    try:
        for ci in range(min(3, len(table.columns))):
            cell = table.cell(row_idx, ci)
            if cell.fill and cell.fill.fore_color and cell.fill.fore_color.rgb:
                rgb = str(cell.fill.fore_color.rgb)
                if rgb.upper() in ("BDD7EE", "4472C4", "5B9BD5", "9DC3E6"):
                    return True
    except Exception:
        pass

    return False


# ═══════════════════════════════════════════════════════════════
# STYLE APPLIERS
# ═══════════════════════════════════════════════════════════════

def apply_cell_font_style(cell, style_cfg, is_header=False):
    """
    Apply font styling to all runs in a cell.
    Preserves strikethrough if already set.
    """
    font_name = style_cfg.get("font_name")
    font_size = style_cfg.get("font_size_pt")
    font_bold = style_cfg.get("font_bold")
    font_italic = style_cfg.get("font_italic")
    font_color = style_cfg.get("font_color")
    alignment = style_cfg.get("alignment")
    v_alignment = style_cfg.get("vertical_alignment")
    bg_color = style_cfg.get("bg_color")

    for p in cell.text_frame.paragraphs:
        # Set paragraph alignment
        if alignment:
            p.alignment = get_alignment(alignment)

        for run in p.runs:
            # Preserve existing strikethrough
            has_strike = run.font._element.attrib.get('strike', None)

            if font_name:
                run.font.name = font_name
            if font_size is not None:
                run.font.size = Pt(font_size)
            if font_bold is not None:
                run.font.bold = font_bold
            if font_italic is not None:
                run.font.italic = font_italic
            if font_color:
                run.font.color.rgb = hex_to_rgb(font_color)

            # Restore strikethrough if it was set
            if has_strike:
                run.font._element.attrib['strike'] = has_strike

    # Vertical alignment
    if v_alignment:
        cell.vertical_anchor = get_vertical_alignment(v_alignment)

    # Background color (header only typically)
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(bg_color)


def apply_textbox_font_style(shape, style_cfg):
    """Apply font styling to all runs in a text box shape."""
    font_name = style_cfg.get("font_name")
    font_size = style_cfg.get("font_size_pt")
    font_bold = style_cfg.get("font_bold")
    font_italic = style_cfg.get("font_italic")
    font_color = style_cfg.get("font_color")

    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            if font_name:
                run.font.name = font_name
            if font_size is not None:
                run.font.size = Pt(font_size)
            if font_bold is not None:
                run.font.bold = font_bold
            if font_italic is not None:
                run.font.italic = font_italic
            if font_color:
                run.font.color.rgb = hex_to_rgb(font_color)


# ═══════════════════════════════════════════════════════════════
# MAIN FIXER
# ═══════════════════════════════════════════════════════════════

def fix_ppt_fonts(input_path, config=None):
    """
    Open a generated PPT, fix all fonts according to config, save as -styled.pptx.
    """
    cfg = config or STYLE_CONFIG
    prs = Presentation(input_path)
    total_slides = len(prs.slides)

    skip_slides = set(cfg.get("skip_slides", []))
    skip_pic_only = cfg.get("skip_picture_only_slides", True)
    skip_placeholders = cfg.get("skip_placeholder_slides", True)

    header_cfg = cfg.get("table_header", {})
    body_cfg = cfg.get("table_body", {})
    title_cfg = cfg.get("section_title", {})
    rag_cfg = cfg.get("rag_legend", {})

    stats = {
        "slides_processed": 0,
        "tables_styled": 0,
        "header_rows_fixed": 0,
        "body_rows_fixed": 0,
        "titles_fixed": 0,
    }

    for si, slide in enumerate(prs.slides):
        slide_num = si + 1

        # ── Skip logic ──
        if slide_num in skip_slides:
            print(f"  Slide {slide_num}/{total_slides}: skipped (config)")
            continue

        if skip_pic_only and len(slide.shapes) == 1:
            if slide.shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"  Slide {slide_num}/{total_slides}: skipped (picture-only)")
                continue

        if skip_placeholders:
            has_ph = any(s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER for s in slide.shapes)
            if has_ph:
                print(f"  Slide {slide_num}/{total_slides}: skipped (placeholder/appendix)")
                continue

        print(f"  Slide {slide_num}/{total_slides}: styling...")
        stats["slides_processed"] += 1

        for shape in slide.shapes:
            # ── Style tables ──
            if shape.has_table:
                table = shape.table
                stats["tables_styled"] += 1

                for ri in range(len(table.rows)):
                    if is_header_row(table, ri):
                        # Apply header style
                        for ci in range(len(table.columns)):
                            apply_cell_font_style(table.cell(ri, ci), header_cfg, is_header=True)
                        stats["header_rows_fixed"] += 1
                    else:
                        # Apply body style
                        for ci in range(len(table.columns)):
                            apply_cell_font_style(table.cell(ri, ci), body_cfg, is_header=False)
                        stats["body_rows_fixed"] += 1

            # ── Style section titles ──
            elif is_section_title(shape) and title_cfg:
                apply_textbox_font_style(shape, title_cfg)
                stats["titles_fixed"] += 1

            # ── Style RAG legend ──
            elif is_rag_legend(shape) and rag_cfg:
                apply_textbox_font_style(shape, rag_cfg)

    # ── Save ──
    base, ext = os.path.splitext(input_path)
    output_path = f"{base}-styled{ext}"
    prs.save(output_path)

    print(f"\nDone! Saved: {output_path}")
    print(f"  Slides processed: {stats['slides_processed']}")
    print(f"  Tables styled: {stats['tables_styled']}")
    print(f"  Header rows fixed: {stats['header_rows_fixed']}")
    print(f"  Body rows fixed: {stats['body_rows_fixed']}")
    print(f"  Section titles fixed: {stats['titles_fixed']}")

    return output_path


# ═══════════════════════════════════════════════════════════════
# CONFIG FILE SUPPORT
# ═══════════════════════════════════════════════════════════════

def load_config(config_path):
    """Load style config from a JSON file and merge with defaults."""
    with open(config_path, "r") as f:
        user_cfg = json.load(f)

    # Deep merge: user values override defaults
    merged = json.loads(json.dumps(STYLE_CONFIG))  # deep copy
    for section_key, section_val in user_cfg.items():
        if isinstance(section_val, dict) and section_key in merged and isinstance(merged[section_key], dict):
            merged[section_key].update(section_val)
        else:
            merged[section_key] = section_val

    return merged


def save_default_config(path="ppt_style_config.json"):
    """Export the default config to a JSON file for easy editing."""
    with open(path, "w") as f:
        json.dump(STYLE_CONFIG, f, indent=2)
    print(f"Default config saved to: {path}")
    print("Edit this file and pass it via: python fix_ppt_fonts.py report.pptx --config ppt_style_config.json")


# ═══════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage:")
        print("  python fix_ppt_fonts.py <input.pptx>                     # Use default styles")
        print("  python fix_ppt_fonts.py <input.pptx> --config style.json # Use custom config")
        print("  python fix_ppt_fonts.py --export-config                  # Export default config to JSON")
        sys.exit(1)

    if sys.argv[1] == "--export-config":
        out_path = sys.argv[2] if len(sys.argv) > 2 else "ppt_style_config.json"
        save_default_config(out_path)
        sys.exit(0)

    input_file = sys.argv[1]
    if not os.path.exists(input_file):
        print(f"Error: File not found: {input_file}")
        sys.exit(1)

    # Check for --config flag
    config = None
    if "--config" in sys.argv:
        idx = sys.argv.index("--config")
        if idx + 1 < len(sys.argv):
            config_path = sys.argv[idx + 1]
            if not os.path.exists(config_path):
                print(f"Error: Config file not found: {config_path}")
                sys.exit(1)
            config = load_config(config_path)
            print(f"Loaded config from: {config_path}")

    fix_ppt_fonts(input_file, config)
