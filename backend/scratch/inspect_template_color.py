from pptx import Presentation
from pptx.dml.color import MSO_COLOR_TYPE

def inspect_slide_0(template_path):
    prs = Presentation(template_path)
    slide = prs.slides[-1]
    print(f"Shapes on Last Slide (index {len(prs.slides)-1}):")
    for shape in slide.shapes:
        print(f" - {shape.name} (Type: {shape.shape_type})")
        if shape.has_table:
            table = shape.table
            print(f"   Table detected. Rows: {len(table.rows)}, Cols: {len(table.columns)}")
            for r in range(len(table.rows)):
                for c in range(len(table.columns)):
                    cell = table.cell(r, c)
                    fill = cell.fill
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    # Check for fill
                    fill = cell.fill
                    try:
                        color = fill.fore_color
                        print(f"     Cell({r},{c}) Fill Type: {fill.type}")
                        if color.type == 1: # RGB
                            print(f"     Cell({r},{c}) Color (RGB): {color.rgb}")
                        elif color.type == 2: # Theme
                            print(f"     Cell({r},{c}) Theme Color: {color.theme_color}")
                            # Try to get brightness/tint
                            try:
                                print(f"     Cell({r},{c}) Tint: {color.tint}")
                            except: pass
                    except Exception as e:
                        print(f"     Error getting fill: {e}")
                    
                    print(f"     Cell({r},{c}) XML: {tcPr.xml}")
                    
                    # Check text formatting
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            try:
                                print(f"       Run Text: '{run.text}', Font Color: {run.font.color.rgb if run.font.color.type == 1 else run.font.color.theme_color}")
                            except Exception:
                                pass
                    print(f"     Cell({r},{c}) Text: '{cell.text.strip()}'")

if __name__ == '__main__':
    import os
    template_path = "templates/default_template.pptx"
    if not os.path.exists(template_path):
        print(f"Template not found at {template_path}")
        exit()
    
    inspect_slide_0(template_path)
