"""
ppt_sanitizer.py — Comprehensive PPT layout and style fixer.

Accepts an optional style_config dict for user-configurable colors, fonts, and sizes.
Call sanitize_ppt(path) or sanitize_presentation(prs, style_config) after ANY modification.
"""

import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ══════════════════════════════════════════════
# DEFAULT STYLE (overridden by style_config)
# ══════════════════════════════════════════════
DEFAULT_STYLE = {
    "title_font": "Tahoma",
    "title_size": 20,
    "title_bold": True,
    "title_color": "14366B",

    "hdr_font": "Calibri",
    "hdr_size": 11,
    "hdr_bold": True,
    "hdr_text_color": "000000",
    "hdr_bg_color": "b8cce4",

    "data_font": "Calibri",
    "data_size": 11,
    "data_size_min": 9,
    "data_size_max": 13,
}

# Layout constants (not user-configurable)
MIN_DATA_ROW_HEIGHT = int(Inches(0.26))
HDR_MIN_HEIGHT = int(Inches(0.32))
HDR_CELL_PAD_EMU = int(Inches(0.10))
MIN_GAP = int(Inches(0.08))
TITLE_TO_TABLE_GAP = int(Inches(0.04))
TABLE_TO_TITLE_GAP = int(Inches(0.22))
NONE_GAP = int(Inches(0.06))
TOP_MARGIN = int(Inches(0.55))
RAG_CLEARANCE = int(Inches(0.10))
SLIDE_BOTTOM_PAD = int(Inches(0.10))

RAG_NAMES = {"Flowchart: Connector 13", "Flowchart: Connector 3", "Flowchart: Connector 5"}
RAG_MARKER = "Project at risk"


def _merge_style(user_config):
    """Merge user config over defaults."""
    cfg = dict(DEFAULT_STYLE)
    if user_config:
        for k, v in user_config.items():
            if k in cfg and v is not None and v != "":
                cfg[k] = v
    return cfg


def _hex_to_rgb(hex_str):
    """Convert '14366B' to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ══════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════
def _in(emu_val):
    return emu_val / 914400

def _emu(inches):
    return int(Inches(inches))

def _is_rag(shape):
    if shape.name in RAG_NAMES:
        return True
    if shape.has_text_frame and RAG_MARKER in shape.text_frame.text:
        return True
    return False

def _is_skip_slide(slide, idx):
    # Skip Title Slide (0) and Summary Slide (1) which has a custom layout
    if idx < 2:
        return True
    shapes = list(slide.shapes)
    if len(shapes) <= 1 and shapes and shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if any(s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER for s in shapes):
        return True
    return False

def _is_section_title(shape):
    if not shape.has_text_frame or shape.has_table:
        return False
    if _is_rag(shape):
        return False
    text = shape.text_frame.text.strip()
    if not text or len(text) > 80:
        return False
    for skip in ["Project at risk", "recover/mitigate", "Project on track",
                  "*See Milestones", "\u00a9", "Vision to Value", "Red:", "Amber:", "Green:"]:
        if skip in text:
            return False
    if "\u2022 None" in text or "None" == text:
        return False
    return True

def _is_none_text(shape):
    if not shape.has_text_frame:
        return False
    t = shape.text_frame.text.strip()
    return t in ("\u2022 None", "\u2022None", "None", "- None")

def _is_content(shape):
    return shape.has_table or _is_section_title(shape) or _is_none_text(shape)


def _is_financial_table(table):
    """
    Returns True if the table is a financial summary (FS or Phase) table.
    Detected by the first row containing SoW/financial/phase-related keywords.
    These tables use fully custom styling and must NOT be touched by the sanitizer.
    """
    if not table.rows:
        return False
    try:
        first_text = table.rows[0].cells[0].text.strip().lower()
    except Exception:
        return False
    financial_keywords = [
        "sow", "financial summary", "financial",
        "design", "build", "testing", "go-live", "golive",
        "hypercare", "hyper care", "sprint", "sow-", "sow ",
        "master data", "data feed", "uат", "uat", "sit"
    ]
    return any(kw in first_text for kw in financial_keywords)

def _actual_height(shape):
    if shape.has_table:
        return sum(shape.table.rows[r].height for r in range(len(shape.table.rows)))
    return shape.height

def _actual_bottom(shape):
    return shape.top + _actual_height(shape)

def _content_shapes(slide):
    return sorted([s for s in slide.shapes if _is_content(s)], key=lambda s: s.top)


# ══════════════════════════════════════════════
# HEADER HEIGHT CALCULATION
# ══════════════════════════════════════════════
def _calc_header_lines(text, col_width_emu, hdr_size_pt):
    if not text:
        return 1
    chars_per_inch = max(4, 6.5 * (11 / hdr_size_pt))
    usable_inches = max(0.3, _in(col_width_emu) - 0.10)
    chars_per_line = max(3, int(usable_inches * chars_per_inch))
    words = text.split()
    lines = 1
    cur = 0
    for w in words:
        wl = len(w)
        if cur == 0:
            cur = wl
        elif cur + 1 + wl <= chars_per_line:
            cur += 1 + wl
        else:
            lines += 1
            cur = wl
    return lines

def _calc_header_row_height(table, hdr_size_pt):
    ncols = len(table.columns)
    max_lines = 1
    for ci in range(ncols):
        text = table.cell(0, ci).text.strip()
        col_w = table.columns[ci].width
        lines = _calc_header_lines(text, col_w, hdr_size_pt)
        if lines > max_lines:
            max_lines = lines
    line_h = int(hdr_size_pt * 1.35 * 12700)
    return max(HDR_MIN_HEIGHT, max_lines * line_h + HDR_CELL_PAD_EMU)


# ══════════════════════════════════════════════
# FIXES
# ══════════════════════════════════════════════
def _sync_table_heights(slide):
    count = 0
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        actual = sum(shape.table.rows[r].height for r in range(len(shape.table.rows)))
        if abs(shape.height - actual) > _emu(0.02):
            shape.height = actual
            count += 1
    return count


def _set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for old in tcPr.findall(".//a:solidFill", nsmap):
        old.getparent().remove(old)
    fill_xml = (
        '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{hex_color}"/>'
        '</a:solidFill>'
    )
    tcPr.append(parse_xml(fill_xml))


def _fix_header_row(table, cfg):
    from pptx.oxml.ns import qn

    ncols = len(table.columns)
    if len(table.rows) < 1:
        return
    hdr_size_pt = cfg["hdr_size"]
    table.rows[0].height = _calc_header_row_height(table, hdr_size_pt)
    for ci in range(ncols):
        cell = table.cell(0, ci)
        text = cell.text.strip()
        tf = cell.text_frame
        tf.clear()

        col_width = table.columns[ci].width
        is_single_word = len(text.split()) <= 1

        if is_single_word:
            # Single-word headers (RAG, ID, S.No, Status) → NEVER wrap
            tf.word_wrap = False
            txBody = cell._tc.find(qn('a:txBody'))
            if txBody is not None:
                bodyPr = txBody.find(qn('a:bodyPr'))
                if bodyPr is not None:
                    bodyPr.set('wrap', 'none')

            # Shrink font if needed to fit single word in column
            EMU_PER_CHAR_PER_PT = 7500
            margin_total = _emu(0.04) * 2
            usable = col_width - margin_total
            text_w = len(text) * EMU_PER_CHAR_PER_PT * hdr_size_pt
            actual_size = hdr_size_pt
            if text_w > usable and usable > 0:
                actual_size = max(7, int(usable / (len(text) * EMU_PER_CHAR_PER_PT)))

            # Tighter margins for very narrow columns
            if col_width < _emu(0.6):
                cell.margin_left = _emu(0.02)
                cell.margin_right = _emu(0.02)
            else:
                cell.margin_left = _emu(0.04)
                cell.margin_right = _emu(0.04)

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.name = cfg["hdr_font"]
            run.font.size = Pt(actual_size)
            run.font.bold = cfg["hdr_bold"]
            run.font.color.rgb = _hex_to_rgb(cfg["hdr_text_color"])
        else:
            # Multi-word: check if it fits on one line
            EMU_PER_CHAR_PER_PT = 7500
            margin_total = _emu(0.04) * 2
            usable = col_width - margin_total
            text_w = len(text) * EMU_PER_CHAR_PER_PT * hdr_size_pt

            if text_w <= usable:
                tf.word_wrap = False
                txBody = cell._tc.find(qn('a:txBody'))
                if txBody is not None:
                    bodyPr = txBody.find(qn('a:bodyPr'))
                    if bodyPr is not None:
                        bodyPr.set('wrap', 'none')
            else:
                tf.word_wrap = True

            cell.margin_left = _emu(0.04)
            cell.margin_right = _emu(0.04)

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.name = cfg["hdr_font"]
            run.font.size = Pt(hdr_size_pt)
            run.font.bold = cfg["hdr_bold"]
            run.font.color.rgb = _hex_to_rgb(cfg["hdr_text_color"])

        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        cell.margin_top = _emu(0.03)
        cell.margin_bottom = _emu(0.03)
        _set_cell_bg(cell, cfg["hdr_bg_color"])
        
        # Apply borders to header cell
        tcPr = cell._tc.get_or_add_tcPr()
        for line in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
            # Remove existing borders of this type
            for old in tcPr.findall(f".//{line}", {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}):
                old.getparent().remove(old)
            # Add new black border
            ln_xml = (
                f'<{line} w="12700" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                '  <a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
                '  <a:prstDash val="solid"/>'
                f'</{line}>'
            )
            tcPr.append(parse_xml(ln_xml))


def _fix_data_cells(table, cfg):
    ncols = len(table.columns)
    nrows = len(table.rows)
    for ri in range(1, nrows):
        if table.rows[ri].height < MIN_DATA_ROW_HEIGHT:
            table.rows[ri].height = MIN_DATA_ROW_HEIGHT
        for ci in range(ncols):
            cell = table.cell(ri, ci)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    if not run.font.name or run.font.name.startswith("+"):
                        run.font.name = cfg["data_font"]
                    elif run.font.name not in (cfg["data_font"], "Tahoma"):
                        run.font.name = cfg["data_font"]
                    if run.font.size:
                        pt = run.font.size.pt
                        if pt < cfg["data_size_min"]:
                            run.font.size = Pt(cfg["data_size_min"])
                        elif pt > cfg["data_size_max"]:
                            run.font.size = Pt(cfg["data_size_max"])
                    else:
                        run.font.size = Pt(cfg["data_size"])
                    try:
                        if run.font.color and run.font.color.rgb:
                            pass
                    except:
                        run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Apply borders to data cell
            tcPr = cell._tc.get_or_add_tcPr()
            for line in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
                for old in tcPr.findall(f".//{line}", {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}):
                    old.getparent().remove(old)
                ln_xml = (
                    f'<{line} w="12700" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    '  <a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
                    '  <a:prstDash val="solid"/>'
                    f'</{line}>'
                )
                tcPr.append(parse_xml(ln_xml))


def _fix_section_titles(slide, cfg):
    for shape in slide.shapes:
        if not _is_section_title(shape):
            continue
        tf = shape.text_frame
        text = tf.text.strip()
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = cfg["title_font"]
        run.font.size = Pt(cfg["title_size"])
        run.font.bold = cfg["title_bold"]
        run.font.color.rgb = _hex_to_rgb(cfg["title_color"])


def _restack(slide, slide_h):
    content = _content_shapes(slide)
    if not content:
        return
    rag_bot = max((s.top + s.height for s in slide.shapes if _is_rag(s)), default=0)
    start_y = max(TOP_MARGIN, rag_bot + RAG_CLEARANCE) if rag_bot > 0 else TOP_MARGIN
    max_bot = slide_h - SLIDE_BOTTOM_PAD
    positions = []
    cursor = start_y
    for i, shape in enumerate(content):
        if i > 0:
            prev = content[i - 1]
            if _is_section_title(prev) and shape.has_table:
                gap = TITLE_TO_TABLE_GAP
            elif _is_section_title(prev) and _is_none_text(shape):
                gap = NONE_GAP
            elif (prev.has_table or _is_none_text(prev)) and _is_section_title(shape):
                gap = TABLE_TO_TITLE_GAP
            else:
                gap = MIN_GAP
            cursor += gap
        h = _actual_height(shape)
        positions.append({"top": cursor, "h": h, "shape": shape})
        cursor += h

    if cursor > max_bot:
        overflow = cursor - max_bot
        gaps = []
        for i in range(1, len(positions)):
            g = positions[i]["top"] - (positions[i-1]["top"] + positions[i-1]["h"])
            gaps.append(g)
        total_gap = sum(gaps)
        min_each = _emu(0.01)
        compressible = total_gap - min_each * len(gaps)
        if compressible > 0:
            ratio = min(1.0, overflow / compressible)
            cursor = positions[0]["top"]
            for i, pos in enumerate(positions):
                pos["top"] = cursor
                cursor += pos["h"]
                if i < len(gaps):
                    red = int((gaps[i] - min_each) * ratio)
                    cursor += max(min_each, gaps[i] - red)

    for pos in positions:
        pos["shape"].top = pos["top"]
        if pos["shape"].has_table:
            pos["shape"].height = pos["h"]


def _fix_rag_collision(slide):
    rags = [s for s in slide.shapes if _is_rag(s)]
    if not rags:
        return
    rag_bot = max(s.top + s.height for s in rags)
    content = _content_shapes(slide)
    if content and rag_bot > content[0].top:
        shift = rag_bot - content[0].top + RAG_CLEARANCE
        for s in content:
            s.top += shift


# ══════════════════════════════════════════════
# PUBLIC API
# ══════════════════════════════════════════════
def sanitize_presentation(prs, style_config=None):
    """Sanitize a Presentation in-place. style_config overrides default colors/fonts."""
    cfg = _merge_style(style_config)
    slide_h = prs.slide_height
    total = len(prs.slides)
    for si, slide in enumerate(prs.slides):
        if _is_skip_slide(slide, si):
            continue
        sn = si + 1
        print(f"  Sanitizing slide {sn}/{total}...")
        
        # Use #00b0f0 for slide 2 headers, default for others
        slide_cfg = dict(cfg)
        if sn == 2:
            slide_cfg["hdr_bg_color"] = "00b0f0"
            
        for shape in slide.shapes:
            if shape.has_table:
                if _is_financial_table(shape.table):
                    # Financial tables (FS/Phase) have fully custom styling — skip sanitizer
                    continue
                _fix_header_row(shape.table, slide_cfg)
                _fix_data_cells(shape.table, slide_cfg)
        _fix_section_titles(slide, slide_cfg)
        synced = _sync_table_heights(slide)
        if synced:
            print(f"    synced {synced} table heights")
        _fix_rag_collision(slide)
        _restack(slide, slide_h)
        _verify(slide, sn, slide_h)
    return prs


def sanitize_ppt(input_path, output_path=None, style_config=None):
    """Open, sanitize, save."""
    if output_path is None:
        output_path = input_path
    prs = Presentation(input_path)
    sanitize_presentation(prs, style_config)
    prs.save(output_path)
    print(f"Sanitized: {output_path}")
    return output_path


def _verify(slide, num, slide_h):
    content = _content_shapes(slide)
    for s in content:
        bot = _actual_bottom(s)
        if bot > slide_h:
            name = s.text_frame.text[:25] if s.has_text_frame else s.name
            print(f"    warning: '{name}' overflow by {_in(bot - slide_h):.2f}in")
    for i in range(len(content) - 1):
        a_bot = _actual_bottom(content[i])
        b_top = content[i + 1].top
        if a_bot > b_top + _emu(0.02):
            an = content[i].text_frame.text[:20] if content[i].has_text_frame else content[i].name
            bn = content[i+1].text_frame.text[:20] if content[i+1].has_text_frame else content[i+1].name
            print(f"    warning: overlap '{an}' <-> '{bn}'")


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python ppt_sanitizer.py <input.pptx> [output.pptx]")
        sys.exit(1)
    sanitize_ppt(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
