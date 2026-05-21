from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import pandas as pd
import datetime
from datetime import datetime as dt
from copy import deepcopy
from dateutil.parser import parse
from pptx.dml.color import MSO_COLOR_TYPE
import re

# ----------------- CONSTANTS -----------------
SLIDE_TOTAL_WIDTH = 10  # Adjust based on your actual PPT slide size in inches
LEFT_MARGIN = 0.7
RIGHT_MARGIN = 0.4
MAX_TABLE_WIDTH = SLIDE_TOTAL_WIDTH - (LEFT_MARGIN + RIGHT_MARGIN)
SLIDE_HEIGHT = 7.5
TOP_MARGIN = 1.2  # Space below title
FOOTER_MARGIN = 0.7  # Space above footer text
MAX_TABLE_HEIGHT = SLIDE_HEIGHT - TOP_MARGIN - FOOTER_MARGIN


# ----------------- HELPERS -----------------

def is_date_string(s):
    """
    Checks if the string s can be parsed as a date.
    Returns True if parse succeeds, False otherwise.
    """
    try:
        parse(s, fuzzy=False)
        return True
    except Exception:
        return False

def is_single_word(col_name, data_rows):
    """
    Returns True if the header and all values in the column are single words (no spaces)
    and reasonably short (<=12 characters). False otherwise.
    """
    header = str(col_name).strip()
    if " " in header or len(header) > 12:
        return False

    for row in data_rows:
        val = row.get(col_name, "")
        # Support dict with "value" key or plain value
        text = val.get("value") if isinstance(val, dict) else val
        text = str(text).strip()

        if not text:
            continue  # Empty cells skip

        # If there's a space or length too long, treat as not single word
        if " " in text or len(text) > 12:
            return False

    return True


def detect_date_columns(columns, data_rows, threshold=0.3):
    """
    Detect columns where at least 'threshold' fraction of non-empty values are date strings.
    Returns list of date columns.
    """
    date_cols = []
    for c in columns:
        total_vals = 0
        date_vals = 0
        if "date" in str(c).lower():
            date_cols.append(c)
            continue
        for row in data_rows:
            val = row.get(c, "")
            text = val.get("value") if isinstance(val, dict) else str(val)
            text = str(text).strip()
            if text:
                total_vals += 1
                if is_date_string(text):
                    date_vals += 1
        if total_vals > 0 and (date_vals / total_vals) >= threshold:
            date_cols.append(c)
    return date_cols


def set_cell_border(cell, border_color="000000", border_width="12700"):
    """
    Sets border for the given pptx cell with specified color and width.
    Existing border elements for the cell are removed before adding new ones.
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove duplicated border nodes
    for tag in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        for element in list(tcPr):
            if element.tag.endswith(tag):
                tcPr.remove(element)
    for line in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = parse_xml(
            f'''<{line} w="{border_width}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                  <a:solidFill><a:srgbClr val="{border_color}"/></a:solidFill>
                  <a:prstDash val="solid"/>
               </{line}>'''
        )
        tcPr.append(ln)


def clear_table_style(table):
    """
    Aggressively clears template-level table styles and fills to ensure manual cell-level
    styling is visible and not overridden by blue/black theme defaults.
    """
    try:
        from pptx.oxml.ns import qn
        tblPr = table._tbl.find(qn('a:tblPr'))
        if tblPr is not None:
            # Set style to 'No Style, No Grid'
            sid = tblPr.find(qn('a:tableStyleId'))
            if sid is not None:
                sid.text = '{00000000-0000-0000-0000-000000000000}'
            # Disable all automated banding/headers/footers
            for flag in ('firstRow', 'firstCol', 'lastRow', 'lastCol', 'bandRow', 'bandCol'):
                tblPr.set(flag, '0')
            # Remove any table-level solid/gradient fills that mask cells
            for tag in ('a:solidFill', 'a:gradFill', 'a:pattFill', 'a:noFill'):
                for el in tblPr.findall(qn(tag)):
                    tblPr.remove(el)
    except Exception as e:
        print(f"Warning: Could not clear table style: {e}")

def set_cell_background(cell, hex_color):
    """
    Sets cell background color using OXML to bypass stubborn theme defaults.
    hex_color: string like 'FFFFFF' or '00AEEF'
    """
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        # Remove existing fills
        for tag in ('a:solidFill', 'a:gradFill', 'a:pattFill', 'a:noFill'):
            for el in tcPr.findall(qn(tag)):
                tcPr.remove(el)
        # Add new solid fill
        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', hex_color)
        solidFill.append(srgbClr)
        tcPr.append(solidFill)
    except Exception as e:
        print(f"Warning: Could not set cell background: {e}")

def clear_table_borders(table):
    """
    Sets all borders in the table to be transparent/invisible to remove black lines.
    """
    try:
        from pptx.oxml.ns import qn
        for row in table.rows:
            for cell in row.cells:
                tcPr = cell._tc.get_or_add_tcPr()
                # Remove all side borders
                for side in ['lnL', 'lnR', 'lnT', 'lnB', 'lnTlToBr', 'lnBlToTr']:
                    ln = tcPr.find(qn(f'a:{side}'))
                    if ln is None:
                        ln = OxmlElement(f'a:{side}')
                        tcPr.append(ln)
                    for child in list(ln):
                        ln.remove(child)
                    # Set to no fill
                    noFill = OxmlElement('a:noFill')
                    ln.append(noFill)
    except Exception as e:
        print(f"Warning: Could not clear table borders: {e}")

def update_footer_year(prs, old_year="2025", new_year="2026"):
    """
    Updates the copyright year in the footer of all slides.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text
                if "2025" in txt and "Prolifics" in txt:
                    # Replace 2025 with 2026
                    shape.text = txt.replace("2025", "2026")

def style_header_cell(cell, text, font_size=10, font_color="FFFFFF", col_width_emu=None, sheet_name=None):
    """
    Standardizes header cell styling:
    - Calibri, centered, bold
    - Specific font size (default 10)
    - Blue background (default Prolifics Blue)
    - White text (default)
    - Yellow background if financial 'Utilised' column
    """
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_VERTICAL_ANCHOR
    
    if sheet_name and sheet_name.lower() == "phase":
        font_size = 8

    is_financial = sheet_name and sheet_name.lower() in ["fs", "phase"]
    is_yellow_header = is_financial and any(k in str(text).lower() for k in ["utilised", "utilize"])

    text_frame = cell.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    
    words = str(text).split()
    is_single_word_val = len(words) <= 1
    
    # Logic to decide if we should wrap
    EMU_PER_CHAR_PER_PT = 3600  # rough estimate for Calibri
    estimated_width = len(str(text)) * font_size * EMU_PER_CHAR_PER_PT
    usable_width = (col_width_emu - Inches(0.1)) if col_width_emu else None
    fits_one_line = usable_width and estimated_width < usable_width

    def _set_no_wrap(c):
        c.text_frame.word_wrap = False
        txBody = c._tc.find(qn('a:txBody'))
        if txBody is not None:
            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                bodyPr.set('wrap', 'none')

    if is_single_word_val:
        if fits_one_line:
            _set_no_wrap(cell)
        actual_size = font_size
        if usable_width and not fits_one_line:
            min_size = 8 if sheet_name and sheet_name.lower() == "phase" else 7
            actual_size = max(min_size, int(usable_width / (len(str(text)) * EMU_PER_CHAR_PER_PT)))
    else:
        if fits_one_line:
            _set_no_wrap(cell)
        else:
            text_frame.word_wrap = True

    is_financial = sheet_name and sheet_name.lower() in ["fs", "phase"]
    is_fs = sheet_name and sheet_name.lower() == "fs"
    h_low = str(text).lower()
    
    # Financial color categories
    is_utilised = any(k in h_low for k in ["utilised", "utilize", "used", "cost"])
    is_remaining = any(k in h_low for k in ["remaining", "balance"])
    is_sow = any(k in h_low for k in ["sow", "charges", "rate", "days"])
    is_services = "services" in h_low
    
    cell.fill.solid()
    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    
    if is_financial:
        # Yellow highlight for "Utilised" header in Summary table
        if is_fs and is_utilised:
            cell.fill.fore_color.rgb = RGBColor(255, 255, 0) # FFFF00 - Pure Yellow
            font_rgb = RGBColor(0, 0, 0) # Black
        elif is_fs and (any(k in h_low for k in ["phase", "po#"]) or is_remaining or is_sow):
            # Summary table headers are mostly white except for Utilised
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) # White
            font_rgb = RGBColor(0, 0, 0) # Black
        elif is_utilised:
            cell.fill.fore_color.rgb = RGBColor(255, 242, 204) # FFF2CC - Tan/Yellow
            font_rgb = RGBColor(0, 0, 0) # Black
        elif is_remaining:
            cell.fill.fore_color.rgb = RGBColor(226, 239, 218) # E2EFDA - Light Green
            font_rgb = RGBColor(0, 0, 0) # Black
        elif is_sow or is_services:
            cell.fill.fore_color.rgb = RGBColor(221, 235, 247) # DDEBF7 - Light Blue
            font_rgb = RGBColor(0, 0, 0) # Black
        else:
            cell.fill.fore_color.rgb = RGBColor(0, 174, 239) # Prolifics Blue
            font_rgb = RGBColor(255, 255, 255) # White
    else:
        cell.fill.fore_color.rgb = RGBColor(0, 174, 239) # Prolifics Blue
        font_rgb = RGBColor(255, 255, 255) # White

    is_phase = sheet_name and sheet_name.lower() == "phase"
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.italic = is_phase  # Phase col headers are italic per design
            r.font.name = "Calibri"
            r.font.size = Pt(font_size)
            r.font.color.rgb = font_rgb
                
    set_cell_border(cell, border_color="00AEEF")


def format_value(value, sheet_name=None, header=None):
    """
    Converts a value to a nicely formatted string, formatting dates as DD-MM-YYYY.
    Returns empty string for NaN values.
    """
    if pd.isna(value):
        return ""

    is_financial = sheet_name and sheet_name.lower() in ["fs", "phase"]
    
    if isinstance(value, (datetime.datetime, datetime.date, pd.Timestamp)):
        if isinstance(value, datetime.datetime):
            value = value.date()
        date_str = value.strftime("%d-%m-%Y")
        return date_str

    elif isinstance(value, str):
        # Try parsing as date
        for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%d", "%d-%m-%Y", "%d/%m/%Y"):
            try:
                parsed = datetime.datetime.strptime(value.strip(), fmt)
                return parsed.strftime("%d-%m-%Y")
            except ValueError:
                continue
        
        # If financial sheet and looks like a number, try to format as currency
        if is_financial:
            try:
                clean_val = value.replace(",", "").replace("£", "").replace("$", "").strip()
                num_val = float(clean_val)
                return f"£{num_val:,.2f}"
            except:
                pass
        return value

    # If numeric and financial, format as currency or float with 2 decimals
    if is_financial and isinstance(value, (int, float)):
        # FORCE currency for financial numbers unless it's clearly just a day count
        h_l = str(header or "").lower()
        if any(k in h_l for k in ["rate", "cost", "charges", "balance", "sow", "remaining", "utilised", "utilize", "total"]):
            return f"£{float(value):,.2f}"
        return f"{float(value):.2f}"
    
    # Try to parse string numbers in financial sheets
    if is_financial and isinstance(value, str):
        try:
            clean_val = value.replace(",", "").replace("£", "").replace("$", "").strip()
            num_val = float(clean_val)
            h_l = str(header or "").lower()
            if any(k in h_l for k in ["rate", "cost", "charges", "balance", "sow", "remaining", "utilised", "utilize", "total"]):
                return f"£{num_val:,.2f}"
            return f"{num_val:.2f}"
        except:
            pass

    return str(value)


# DATE_FORMAT_REGION = "UK"  # default

# def format_value(value):
#     """Format values including date according to selected region (US/UK)."""
#     global DATE_FORMAT_REGION

#     def get_date_format():
#         if DATE_FORMAT_REGION == "US":
#             return "%m-%d-%Y"   # US
#         elif DATE_FORMAT_REGION == "UK":
#             return "%Y%m%d"     # UK (default)
#         return "%Y%m%d"

#     if isinstance(value, (datetime.datetime, datetime.date, pd.Timestamp)):
#         if isinstance(value, datetime.datetime):
#             value = value.date()
#         return value.strftime(get_date_format())

#     elif isinstance(value, str):
#         for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%d", "%d-%m-%Y", "%d/%m/%Y", "%m-%d-%Y"):
#             try:
#                 parsed = datetime.datetime.strptime(value.strip(), fmt)
#                 return parsed.strftime(get_date_format())
#             except ValueError:
#                 continue
#         return value

#     return str(value)

    
def parse_rgb(color, fallback=RGBColor(0, 0, 0)):
    """
    Convert a color from multiple formats to RGBColor:
    - Hex string: "FF000000", "#RRGGBB", "RRGGBB"
    - List/tuple: [R, G, B] or (R, G, B)
    - Returns fallback on error
    """
    if not color:
        return fallback

    # Already a list or tuple [r, g, b]
    if isinstance(color, (list, tuple)) and len(color) == 3:
        try:
            r, g, b = [int(c) for c in color]
            return RGBColor(r, g, b)
        except Exception:
            return fallback

    # Otherwise, assume string
    if isinstance(color, str):
        c = color.replace("#", "").upper()
        if len(c) == 8:  # ARGB
            c = c[2:]
        if len(c) == 6:
            try:
                r = int(c[0:2], 16)
                g = int(c[2:4], 16)
                b = int(c[4:6], 16)
                return RGBColor(r, g, b)
            except Exception:
                return fallback

    # Fallback if all else fails
    return fallback

def style_data_cell(cell, cell_info, font_size=10, sheet_name=None, header=None, is_total_row=False):
    """Styles a data cell, handling currency formatting and total row highlights."""
    if isinstance(cell_info, dict):
        value = cell_info.get("value", "")
        style = cell_info.get("style") or {}
        strike_runs = cell_info.get("strike_runs")
    else:
        value = cell_info
        style = {}
    if sheet_name and sheet_name.lower() == "phase":
        font_size = 8

    is_financial = sheet_name and sheet_name.lower() in ["fs", "phase"]
    h_lower = str(header or "").lower()
    is_yellow_cell = is_total_row and any(k in h_lower for k in ["utilised", "utilize", "remaining", "balance"])

    v_upper = str(value).upper().strip()
    is_rag = v_upper in ["GREEN", "AMBER", "RED", "G", "A", "R", "CLOSE"]

    text_frame = cell.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if (is_financial or is_rag) else PP_ALIGN.LEFT

    default_bold = bool(style.get("font_bold", False)) or is_total_row
    _is_phase_sheet = sheet_name and sheet_name.lower() == "phase"
    default_italic = bool(style.get("font_italic", False)) or (is_total_row and _is_phase_sheet)

    # Apply currency formatting for financial sheets
    formatted_text = format_value(value, sheet_name=sheet_name, header=header)

    # RAG Status Coloring Logic
    # v_upper = str(value).upper().strip() (already defined above)
    bg_color = None
    text_color = RGBColor(0, 0, 0)
    
    # Phase-specific coloring logic: ONLY color the Total row
    if _is_phase_sheet:
        if is_total_row:
            # Priority 1: Remaining / Balance -> Green
            if any(kw in h_lower for kw in ["remaining", "balance"]):
                bg_color = RGBColor(198, 239, 206) # Light Green
            # Priority 2: Used / Cost -> Yellow (Tan)
            elif any(kw in h_lower for kw in ["used", "cost"]):
                bg_color = RGBColor(255, 242, 204) # Light Yellow (Tan)
            # Priority 3: From SOW (Days / Rate / Charges) -> Blue
            elif any(kw in h_lower for kw in ["days", "rate", "charges"]):
                bg_color = RGBColor(221, 235, 247) # Light Blue
        else:
            # Data rows (middle) should be white
            bg_color = RGBColor(255, 255, 255)
    
    # Financial Summary (FS) Summary sheet logic
    elif sheet_name and sheet_name.lower() == "fs":
        # Default FS data cells to white
        bg_color = RGBColor(255, 255, 255)
        
        # Exception 1: PO# column is always light green
        if h_lower == "po#":
            bg_color = RGBColor(198, 239, 206) # Light Green
        
        # Exception 2: Total row remaining/yellow highlight
        if is_total_row and "remaining" in h_lower:
            bg_color = RGBColor(255, 255, 0) # Pure Yellow

    if v_upper in ["GREEN", "AMBER", "RED", "G", "A", "R", "CLOSE"]:
        if v_upper in ["GREEN", "G", "CLOSE"]:
            bg_color = RGBColor(198, 239, 206) # Light Green
            text_color = RGBColor(0, 97, 0)     # Dark Green text
        elif v_upper in ["RED", "R"]:
            bg_color = RGBColor(255, 199, 206) # Light Red
            text_color = RGBColor(156, 0, 6)    # Dark Red text
        elif v_upper in ["AMBER", "A"]:
            bg_color = RGBColor(255, 235, 156) # Light Amber
            text_color = RGBColor(156, 101, 0)  # Dark Amber text

    # Task: Revised Target Date strike-through for old dates
    if h_lower == "revised target date" and not strike_runs:
        lines = str(formatted_text).split('\n')
        if len(lines) > 1:
            strike_runs = []
            for i, line in enumerate(lines):
                strike_runs.append({
                    "text": line,
                    "strike": (i < len(lines) - 1) # Strike all but the last date
                })

    if strike_runs and isinstance(strike_runs, list):
        for i, run_info in enumerate(strike_runs):
            text = str(run_info.get("text", "")).strip()
            if not text: continue
            if i > 0: p.add_line_break()
            run = p.add_run()
            run.text = text
            run.font.name = "Calibri"
            run.font.size = Pt(font_size)
            run.font.bold = (i == len(strike_runs) - 1) if not is_total_row else True
            run.font.color.rgb = text_color if bg_color else RGBColor(0,0,0)
            if run_info.get("strike"):
                run.font._element.attrib['strike'] = 'sngStrike'
    else:
        run = p.add_run()
        run.text = formatted_text
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.bold = default_bold
        run.font.italic = default_italic
        run.font.color.rgb = text_color if bg_color else RGBColor(0, 0, 0)

    cell.fill.solid()
    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    
    if bg_color:
        cell.fill.fore_color.rgb = bg_color
    elif is_financial:
        h_low = str(header or "").lower()
        is_fs = sheet_name and sheet_name.lower() == "fs"
        is_phase = sheet_name and sheet_name.lower() == "phase"

        if is_fs:
            # FS Summary: data cells are white except PO# (green) and total-remaining (yellow)
            if is_total_row and any(k in h_low for k in ["remaining", "balance"]):
                cell.fill.fore_color.rgb = RGBColor(255, 255, 0)       # Pure Yellow
            elif "po#" in h_low:
                cell.fill.fore_color.rgb = RGBColor(198, 224, 180)     # Light Green for PO#
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)     # White
        elif is_phase:
            # Phase Detailed: colour-coded by column group
            # Priority: remaining > used/cost > days/rate/charges > white
            if any(k in h_low for k in ["remaining", "balance"]):
                cell.fill.fore_color.rgb = RGBColor(226, 239, 218)     # Light Green
            elif any(k in h_low for k in ["used", "cost"]):
                cell.fill.fore_color.rgb = RGBColor(255, 242, 204)     # Tan/Yellow
            elif any(k in h_low for k in ["days", "rate", "charges"]):
                cell.fill.fore_color.rgb = RGBColor(221, 235, 247)     # Light Blue
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)     # White (Services etc.)
        else:
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)         # White for other financial
    elif is_yellow_cell:
        cell.fill.fore_color.rgb = RGBColor(255, 255, 0) # Yellow
    else:
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255) # White
        
    set_cell_border(cell)



def insert_slide(prs, layout, index):
    """
    Inserts a slide at a specified index in the presentation.
    Moves the newly added slide to the target index properly.
    """
    new_slide = prs.slides.add_slide(layout)
    sldIdLst = prs.slides._sldIdLst  # XML <p:sldIdLst>
    new_sldId = sldIdLst[-1]        # the sldId for the newly added slide
    sldIdLst.remove(new_sldId)
    index = max(0, min(index, len(sldIdLst)))
    sldIdLst.insert(index, new_sldId)
    return new_slide

def get_day_suffix(day):
        """Return the appropriate suffix for the day number."""
        if 10 <= day % 100 <= 20:
            return "th"
        else:
            return {1: "st", 2: "nd", 3: "rd"}.get(day % 10, "th")
        
def get_next_friday():
    """
    Returns a Friday date based on the rules:
    - Monday, Tuesday, Wednesday: previous Friday
    - Thursday, Friday, Saturday, Sunday: this week's Friday
    """
    today = datetime.date.today()
    weekday = today.weekday()  # Monday=0 ... Sunday=6

    if weekday <= 2:  # Monday(0), Tuesday(1), Wednesday(2)
        # previous Friday
        days_since_friday = (weekday + 3)  # Mon=3, Tue=4, Wed=5
        target = today - datetime.timedelta(days=days_since_friday)
    else:  # Thursday(3) ... Sunday(6)
        # this week's Friday
        days_until_friday = (4 - weekday) % 7
        target = today + datetime.timedelta(days=days_until_friday)
        
    # Call get_day_suffix here 👇
    day = target.day
    suffix = get_day_suffix(day)

    # Format like '07th November 2025'
    return f"{day:02d}{suffix} {target.strftime('%B %Y')}"


def replace_text_with_format(paragraph, new_text, font_name="Calibri", force_bold=False):
    """
    Replace the text in a paragraph while preserving existing formatting of each run.
    The number of runs and their formatting are maintained, and new text is split accordingly.
    """
    runs = paragraph.runs
    if not runs:
        run = paragraph.add_run(new_text)
        run.font.name = font_name
        if force_bold:
            run.font.bold = True
        return run

    # Flatten all existing text (remove line breaks if any)
    old_text = ''.join(run.text for run in runs)
    if not old_text:
        # No text, just replace whole paragraph
        paragraph.text = new_text
        return paragraph.runs[0]

    # Split new_text proportionally to existing runs
    new_text_len = len(new_text)
    old_text_len = len(old_text)
    text_index = 0

    for i, run in enumerate(runs):
        # Calculate proportional split for this run
        run_len = len(run.text)
        if run_len == 0:
            continue

        # Compute next text slice (same relative size)
        slice_end = text_index + int((run_len / old_text_len) * new_text_len)
        if i == len(runs) - 1:
            # last run takes remaining text
            slice_end = new_text_len

        new_slice = new_text[text_index:slice_end]
        run.text = new_slice

        # Optional font adjustments
        if font_name:
            run.font.name = font_name
        if force_bold:
            run.font.bold = True

        text_index = slice_end

    # If there are fewer runs than text length, append remainder
    if text_index < new_text_len:
        extra_run = paragraph.add_run(new_text[text_index:])
        if font_name:
            extra_run.font.name = font_name
        if force_bold:
            extra_run.font.bold = True

    return paragraph.runs[0]


def copy_shapes_from_slide(src_slide, dest_slide, shape_names):
    """
    Copy specific shapes (by name) from src_slide to dest_slide.
    Deepcopies the XML element and inserts it into destination slide.
    """
    for shape in src_slide.shapes:
        if shape.name in shape_names:
            el = shape.element
            new_el = deepcopy(el)
            dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')


def apply_table_borders(table, border_color="000000", border_width="25400"):
    """
    Apply borders to all cells in a table.
    """
    for row in table.rows:
        for cell in row.cells:
            set_cell_border(cell, border_color, border_width)


def add_project_status_title(slide, project_title, left=0.7, top=0.7, width=5, height=0.5):
    """
    Adds the project title + ' - Project Status' at the top-left of a slide.
    Formats text with Tahoma 20pt, bold and dark blue color.
    """
    import re
    # Strip everything until first alphanumeric, then handle 'SB' prefix
    clean_title = project_title.strip()
    # Match and remove any leading non-alphanumeric junk or SB prefix
    clean_title = re.sub(r'^[^a-zA-Z0-9]*', '', clean_title)
    if clean_title.upper().startswith("SB"):
        clean_title = re.sub(r'^SB\s*[\-\s]*', '', clean_title, flags=re.IGNORECASE)
    
    title_text = f"SB - {clean_title.strip()} - Project Status"
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run = p.add_run()
    run.text = title_text
    run.font.name = "Tahoma"
    run.font.bold = True
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(20, 54, 107)  # Dark blue

    return textbox

def compute_header_width(header, min_col_width, max_col_width):
    words = header.split()
    longest_word = max(len(w) for w in words) if words else 1
    if len(words) == 1:
        # single word -> single line
       # Single word -> must fit completely on one line
        width = max(min_col_width, longest_word * 11000)  # tuned for Calibri 14pt
        return min(width, max_col_width)
    
    # Try distributing words across 2–3 lines
    n_lines = min(3, len(words))  
    # Approx width = max line length (in chars) * factor
    lines = [[] for _ in range(n_lines)]
    # for i, w in enumerate(words):
    #     lines[i % n_lines].append(w)
    # line_lengths = [len(" ".join(line)) for line in lines]
    # max_len = max(line_lengths)
    line_lengths = [0] * n_lines

    for w in words:
        # put word in the line with currently smallest length
        idx = line_lengths.index(min(line_lengths))
        lines[idx].append(w)
        line_lengths[idx] += len(w) + 1  # +1 for space

    # Find max line length (in characters)
    max_len = max(len(" ".join(line)) for line in lines)
    
    width = max(min_col_width, max(longest_word * 11000, max_len * 9000))  # factor for EMU scaling
    return min(width, max_col_width)

def calculate_content_based_widths(columns, data_rows, total_width_emu, fixed_widths_in={}):
    """
    Calculate column widths in EMU based on content length, wrapping rules, and fixed widths.

    Special handling:
    - S.No, ID, RAG → minimum width to fit header on one line (data is tiny).
    - Single-word or date → always single line.
    - Multi-word headers → wrapping adjusted into max 3 lines.
    - Content-heavy columns get proportionally more space.
    """
    from pptx.util import Inches
    import re

    # --- Config ---
    MIN_COL_WIDTH = Inches(0.7)   # general minimum width
    MAX_COL_WIDTH = Inches(3.0)   # max allowed for multi-word columns
    SPECIAL_COLS = {"S.No", "ID", "RAG"}
    DATE_PATTERN = re.compile(r"^\d{1,2}[-/]\d{1,2}([-/]\d{2,4})?$")

    # Convert fixed widths to EMU if provided
    fixed_emu = {k: int(v) for k, v in fixed_widths_in.items()}

    # --- Step 1: Collect stats ---
    col_lengths = {c: len(str(c)) for c in columns}  # header length
    col_longest_word = {c: max(len(w) for w in str(c).split()) for c in columns}
    col_word_counts = {c: len(str(c).split()) for c in columns}

    for row in data_rows:
        for c in columns:
            val = row.get(c, "")
            text = val.get("value") if isinstance(val, dict) else val
            text = "" if text is None else str(text)   # <-- FIX: ensure string
            col_lengths[c] += len(text)

            words = text.split()
            if words:
                col_longest_word[c] = max(col_longest_word[c], max(len(w) for w in words))
                col_word_counts[c] = max(col_word_counts[c], len(words))

    # --- Step 2: Assign preliminary widths ---
    widths = {}
    min_widths = {}

    for c in columns:
        header = str(c)

        if c in fixed_emu:
            widths[c] = fixed_emu[c]
            continue

        if c in SPECIAL_COLS:
            min_widths[c] = int(max(MIN_COL_WIDTH, col_longest_word[c] * 0.18))
        elif col_word_counts[c] == 1 or DATE_PATTERN.match(header):
            min_widths[c] = int(max(MIN_COL_WIDTH, col_longest_word[c] * 0.18))
        else:
            min_widths[c] = compute_header_width(header, MIN_COL_WIDTH, MAX_COL_WIDTH)

    # --- Step 3: Distribute space proportionally ---
    total_min = sum(min_widths.values()) + sum(fixed_emu.values())
    remaining = max(total_width_emu - total_min, 0)

    # Weight = total text length per column
    # weights = {c: col_lengths[c] for c in columns if c not in fixed_emu}
    weights = {c: max(col_lengths[c], len(str(c))) for c in columns if c not in fixed_emu}
    total_weight = sum(weights.values()) or 1

    for c in columns:
        if c in fixed_emu:
            widths[c] = fixed_emu[c]
        else:
            extra = int(remaining * (weights[c] / total_weight))
            widths[c] = min_widths[c] + extra


    return widths

def merge_dependencies_table(table, start_row=1, skip_column="Comments"):
    """
    Merge cells downward in Dependencies table for rows where:
    - Only one column has value per row except skip_column
    - Do NOT merge the skip_column
    - Skip headers and already merged cells
    """
    num_rows = len(table.rows)
    num_cols = len(table.columns)

    # Find index of skip_column
    skip_idx = None
    for c in range(num_cols):
        if table.cell(0, c).text.strip() == skip_column:
            skip_idx = c
            break

    for r in range(start_row, num_rows):
        # Identify all non-blank cells except skip_column
        non_blank_cells = [c for c in range(num_cols) 
                           if c != skip_idx and table.cell(r, c).text.strip() != ""]
        if len(non_blank_cells) != 1:
            continue  # skip row if multiple non-blank cells

        col_idx = non_blank_cells[0]
        merge_start = r

        # Merge downward as long as cells below in the same column are blank
        for rr in range(r + 1, num_rows):
            below_cell = table.cell(rr, col_idx)
            if hasattr(below_cell, 'merged') and below_cell.merged:
                break
            if below_cell.text.strip() != "":
                break

        merge_end = rr - 1
        if merge_end > merge_start:
            try:
                table.cell(merge_start, col_idx).merge(table.cell(merge_end, col_idx))
            except ValueError:
                pass

def update_project_title(slide, new_title, author=None):
    """
    Update project title text in the shape named 'Rectangle 6' on slide 1.
    Also updates author name if provided.
    """
    for shape in slide.shapes:
        if shape.name == 'Rectangle 6' and shape.has_text_frame:
            # Update Title (usually 1st paragraph)
            if shape.text_frame.paragraphs:
                replace_text_with_format(shape.text_frame.paragraphs[0], new_title)
            
            # Update Author (usually 4th paragraph)
            if author:
                for para in shape.text_frame.paragraphs:
                    if "author" in para.text.lower():
                        replace_text_with_format(para, f"Author: {author}")
                        break
            return


def update_slide_date(slide):
    """
    Update the first shape containing "Date:" text with the current Friday date.
    """
    friday_str = get_next_friday()
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if "Date:" in p.text:
                    replace_text_with_format(p, f"Date:  {friday_str}")
                    return

def style_section_title_cell(cell, section_title, style_info=None, sheet_name=None):
    """
    Styles a section title cell using given style_info dictionary.
    Applies fill color, font color, font size, bold, italic, underline,
    text alignment, vertical anchor, and sets the section_title text.
    """
    if style_info is None:
        style_info = {}
    
    # Clear existing text/frame
    tf = cell.text_frame
    tf.clear()
    
    # Set fill color
    s_low = section_title.lower()
    is_fs_summary = ("summary" in s_low and "detailed" not in s_low) or (sheet_name and sheet_name.lower() == "fs")
    is_detailed_phase = "detailed" in s_low or "design" in s_low or "build" in s_low or "sprint" in s_low or (sheet_name and sheet_name.lower() == "phase")
    
    cell.fill.solid()
    if is_fs_summary:
        cell.fill.fore_color.rgb = RGBColor(189, 215, 238) # BDD7EE - Light Blue
        font_rgb = RGBColor(0, 0, 0) # Black
    elif is_detailed_phase:
        cell.fill.fore_color.rgb = RGBColor(198, 239, 206) # C6EFCE - Light Green
        font_rgb = RGBColor(0, 0, 0) # Black
    else:
        fill_color = style_info.get("fill_color")
        cell.fill.fore_color.rgb = parse_rgb(fill_color, fallback=RGBColor(196, 189, 151))
        font_rgb = parse_rgb(style_info.get("font_color"), fallback=RGBColor(0, 0, 0))
    
    # Setup paragraph and run
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if (is_fs_summary or is_detailed_phase) else PP_ALIGN.LEFT
    run = p.add_run()
    run.text = section_title
    
    # Apply font styles
    f_size = style_info.get("font_size", 12)
    if is_detailed_phase:
        f_size = 8
    run.font.name = "Calibri"
    run.font.size = Pt(f_size)
    run.font.bold = True
    run.font.color.rgb = font_rgb

    
    # Horizontal alignment
    horiz = style_info.get("alignment_horizontal", "center").lower()
    if horiz == "center":
        p.alignment = PP_ALIGN.CENTER
    elif horiz == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    
    # Vertical alignment
    vert_map = {
        "top": MSO_VERTICAL_ANCHOR.TOP,
        "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
        "bottom": MSO_VERTICAL_ANCHOR.BOTTOM
    }
    vert = style_info.get("alignment_vertical", "middle").lower()
    cell.vertical_anchor = vert_map.get(vert, MSO_VERTICAL_ANCHOR.MIDDLE)
    
    # Remove paragraph spacing
    p.space_before = 0
    p.space_after = 0
    p.line_spacing = 1
    
    # Set border optionally if needed (you can call your set_cell_border here)
    set_cell_border(cell)
    
def safe_format_date(val):
    """Convert Excel value or text date into DD-MM-YYYY string if possible, else plain string."""
    if val in (None, "", " "):
        return None

    if isinstance(val, dt):
        return val.strftime("%d-%m-%Y")

    if isinstance(val, str):
        for fmt in ("%d %b %Y", "%Y-%m-%d", "%d-%m-%Y", "%d/%m/%Y", "%m/%d/%Y"):
            try:
                return dt.strptime(val.strip(), fmt).strftime("%d-%m-%Y")
            except ValueError:
                continue
    return str(val)  # fallback: return as-is


def get_project_start_date(excel_data):
    """
    Finds the project start date from Key Milestones sheet/table.
    Matches any header containing 'start date', 'startdate', or 'start_date'
    (case-insensitive, partial match allowed).
    """

    for sheet_name, tables in excel_data.items():
        for table in tables:
            # Check if table relates to key milestones
            if "key milestone" in sheet_name.lower() or \
               "key milestone" in table.get("section_title", "").lower():

                # Find start date column (supports partial matches)
                start_date_col = None
                for h in table.get("headers", []):
                    if not h:
                        continue
                    header_lower = h.strip().lower()
                    if any(keyword in header_lower for keyword in ["start date", "startdate", "start_date"]):
                        start_date_col = h
                        break

                if not start_date_col:
                    continue

                # Take the first non-empty valid date
                for row in table.get("rows", []):
                    val = row.get(start_date_col, {}).get("value")
                    formatted = safe_format_date(val)
                    if formatted:
                        return formatted

    return None


def copy_font_style(src_run, dest_run):
    """Copy font style from src_run to dest_run safely."""
    dest_font = dest_run.font
    src_font = src_run.font

    # Basic styles
    dest_font.bold = src_font.bold
    dest_font.italic = src_font.italic
    dest_font.underline = src_font.underline
    dest_font.size = src_font.size
    dest_font.name = src_font.name

    # Font color handling (RGB or Theme-based)
    if src_font.color.type == MSO_COLOR_TYPE.RGB:
        dest_font.color.rgb = src_font.color.rgb
    elif src_font.color.type == MSO_COLOR_TYPE.SCHEME:
        dest_font.color.theme_color = src_font.color.theme_color
    # you can add elif for MSO_COLOR_TYPE.AUTO if needed

def update_key_milestones_table(slide, key_milestones):
    
    for shape in slide.shapes:
        if shape.name == "Table 4" and shape.has_table:
            tbl = shape.table

            if len(tbl.rows) < 1:
                print("Error: Table 4 must have at least 1 row to clone.")
                return

            # Style header row
            for col_idx in range(len(tbl.columns)):
                cell = tbl.cell(0, col_idx)
                style_header_cell(cell, cell.text, font_size=12, col_width_emu=tbl.columns[col_idx].width)

            # Column mapping
            headers = [c.text.strip().lower() for c in tbl.rows[0].cells]
            col_map = {}
            for idx, h in enumerate(headers):
                if "key activities" in h or "milestones" in h:
                    col_map["Milestone"] = idx
                elif "milestone targeted date" in h or "target date" in h:
                    col_map["Revised Target Date"] = idx
                elif "date achieved" in h or "delivered" in h:
                    col_map["Date Delivered"] = idx
                elif "status" in h:
                    col_map["Status"] = idx
                elif "rag" in h:
                    col_map["RAG"] = idx
            
            # Set Slide 2 layout for Table 4
            shape.left = int(5.6 * 914400)
            shape.top = int(3.1 * 914400)
            shape.width = int(7.2 * 914400)
            
            # Compact column widths for Table 4 (total ~7.2)
            if len(tbl.columns) >= 5:
                tbl.columns[0].width = int(Inches(2.2))
                tbl.columns[1].width = int(Inches(1.2))
                tbl.columns[2].width = int(Inches(1.2))
                tbl.columns[3].width = int(Inches(1.8))
                tbl.columns[4].width = int(Inches(0.8))
                # Shrink any extra columns
                for extra_idx in range(5, len(tbl.columns)):
                    tbl.columns[extra_idx].width = 1

            # Clone last row (or first data row if last not exists)
            data_row_idx = 1 if len(tbl.rows) > 1 else 0
            last_row = tbl.rows[data_row_idx]

            # Add enough rows
            rows_needed = len(key_milestones) - (len(tbl.rows) - 1)
            for _ in range(max(0, rows_needed)):
                new_row = deepcopy(last_row._tr)
                tbl._tbl.append(new_row)

            # Fill table
            for row_idx, milestone in enumerate(key_milestones, start=1):
                # Apply white background (with light blue banding) and blue borders
                row_bg = "FFFFFF" if row_idx % 2 == 1 else "D9F3FD"
                for ci in range(len(tbl.columns)):
                    c = tbl.cell(row_idx, ci)
                    set_cell_background(c, row_bg)
                    set_cell_border(c, border_color="00AEEF", border_width="19050") # Thicker border
                    c.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

                for key, col_idx in col_map.items():
                    cell = tbl.cell(row_idx, col_idx)
                    # Use smaller font for summary table to avoid overflow
                    font_size = 10
                    if key == "Revised Target Date":
                        revised = milestone.get("Revised Target Date", {}).get("value", "")
                        target  = milestone.get("Target Date", {}).get("value", "")
                        enddate = milestone.get("End Date", {}).get("value", "")

                        val = revised or target or enddate
                        val = format_value(val)
                        
                        val = revised or target or enddate
                        val = format_value(val)
                        
                        # Only show multiple dates (with strikethrough) for Go-Live row
                        # For all other milestones, just show the latest date
                        milestone_name = str(milestone.get("Milestone", {}).get("value", "")).lower()
                        if "go-live" not in milestone_name and "go live" not in milestone_name:
                            if "\n" in str(val):
                                val = str(val).strip().split("\n")[-1].strip()

                    elif key == "Date Delivered":
                        delivered = milestone.get("Date Delivered", {}).get("value", "")
                        val = format_value(delivered)
                    else:
                        val = milestone.get(key, {}).get("value", "")

                    # else:
                    #     val = milestone.get(key, {}).get("value", "")
                    # cell.text = str(val) if val not in (None, "") else ""
                    if key == "RAG":
                        img_map = {
                            "G": "./images/green_icon.png",
                            "R": "./images/red_icon.png",
                            "A": "./images/amber_icon.png"
                        }
                        if val in img_map:
                            cell.text = ""  # clear text

                            # compute cell coordinates
                            x = tbl._graphic_frame.left + sum(tbl.columns[i].width for i in range(col_idx))
                            y = tbl._graphic_frame.top + sum(tbl.rows[j].height for j in range(row_idx))
                            cx = tbl.columns[col_idx].width
                            cy = tbl.rows[row_idx].height
                            
                            # Make image smaller (e.g. 40% of cell size)
                            img_width = cx * 0.4
                            img_height = cy * 0.4

                            # # center the image inside the cell
                            # slide.shapes.add_picture(
                            #     img_map[val],
                            #     x + cx//4,  # offset to center
                            #     y + cy//4,
                            #     width=cx//2,
                            #     height=cy//2
                            # )
                            # Center the smaller image
                            slide.shapes.add_picture(
                                img_map[val],
                                x + (cx - img_width) // 2,
                                y + (cy - img_height) // 2,
                                width=img_width,
                                height=img_height
                            )
                        else:
                            cell.text = str(val) if val else ""
                            
                    # --- All other columns ---
                    else:
                        txt_val = str(val) if val not in (None, "") else ""
                        cell.text = txt_val
                        
                        # Apply strikethrough for multi-line dates (e.g. Go-Live)
                        if "\n" in txt_val:
                            lines = txt_val.split("\n")
                            cell.text_frame.clear()
                            p = cell.text_frame.paragraphs[0]
                            for i, line in enumerate(lines):
                                if i > 0: p.add_line_break()
                                run = p.add_run()
                                run.text = line.strip()
                                if i < len(lines) - 1:
                                    # Add strikethrough to old dates
                                    run.font._element.attrib['strike'] = 'sngStrike'
                    
                    # Set alignment: Milestone column left-aligned, others centered
                    for paragraph in cell.text_frame.paragraphs:
                        if key == "Milestone":
                            paragraph.alignment = PP_ALIGN.LEFT
                        else:
                            paragraph.alignment = PP_ALIGN.CENTER
                        
                        for run in paragraph.runs:
                            run.font.name = "Calibri"
                            run.font.size = Pt(12)
                            run.font.color.rgb = RGBColor(0, 0, 0)  # black
                            
                    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                
                # Force tighter row height to reduce box size
                tbl.rows[row_idx].height = int(Inches(0.35))
            
            # Ensure no template styles interfere (but keep our custom blue borders)
            clear_table_style(tbl)

            # Add the footer note below Table 4 (removing any existing first)
            try:
                note_text = "See Milestones , Issues & Risks sections for further details."
                
                # Broad cleanup of any shapes containing this text
                for s in list(slide.shapes):
                    try:
                        if s.has_text_frame and "See Milestones" in s.text:
                            sp = s._element
                            sp.getparent().remove(sp)
                    except:
                        pass

                left = tbl._graphic_frame.left
                # Explicitly sum row heights and add a generous gap (0.4 inches)
                actual_table_height = sum(r.height for r in tbl.rows)
                top = tbl._graphic_frame.top + actual_table_height + int(Inches(0.4))
                width = tbl._graphic_frame.width
                height = int(Inches(0.3))
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = note_text
                run.font.italic = True
                run.font.size = Pt(11)
                run.font.name = "Calibri"
                run.font.color.rgb = RGBColor(0, 0, 0)
            except Exception as e:
                print(f"Warning: Could not add milestone footer note: {e}")
            break


def update_project_start_date(slide, project_start_date):
    """
    Updates the Project Start Date table (Table 18) with Prolifics Blue header 
    and White data cell.
    """
    for shape in slide.shapes:
        if shape.name.lower().strip() == "table 18" and shape.has_table:
            table = shape.table
            clear_table_style(table)
            # Assume cell(0,0) is label, cell(0,1) is value
            label_cell = table.cell(0, 0)
            value_cell = table.cell(0, 1)

            # Style Header (Label)
            set_cell_background(label_cell, "00AEEF") # Prolifics Blue
            for p in label_cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(12)
                    r.font.color.rgb = RGBColor(255, 255, 255)

            # Style Value
            value_cell.text = str(project_start_date) if project_start_date else ""
            set_cell_background(value_cell, "FFFFFF") # White
            for p in value_cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(14)
                    r.font.color.rgb = RGBColor(0, 0, 0)

            # Style Table 18 (Project Start Date)
            for row in table.rows:
                for cell in row.cells:
                    set_cell_background(cell, "00AEEF") # Blue background
                    set_cell_border(cell, border_color="00AEEF", border_width="12700")
                    cell.margin_left = int(Inches(0.05))
                    # Make text white if it's a label
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.color.rgb = RGBColor(255, 255, 255)
                            r.font.bold = True
            
            # Set Slide 2 layout for Table 18 to match Table 4 width
            shape.left = int(5.6 * 914400)
            shape.top = int(1.4 * 914400)
            shape.width = int(7.2 * 914400)
            
            # Value cell should be white background with black text
            if len(table.columns) > 1:
                val_cell = table.cell(0, 1)
                set_cell_background(val_cell, "FFFFFF")
                # Balanced column widths for 7.2 total width
                table.columns[0].width = int(Inches(3.0))
                table.columns[1].width = int(Inches(4.2))
                for p in val_cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = RGBColor(0, 0, 0)
                        r.font.bold = True
            return

# --- Helper functions ---

def filter_valid_data_rows(chunk, columns):
    """Return rows with at least one non-empty value (ignoring section_title-only rows)"""
    valid_rows = []
    for r in chunk:
        if "section_title" in r:
            continue
        if any(
            (isinstance(r.get(c), dict) and str(r.get(c).get("value", "")).strip()) or
            (not isinstance(r.get(c), dict) and str(r.get(c)).strip() != "")
            for c in columns if not c.lower().startswith("unnamed")
        ):
            valid_rows.append(r)
    return valid_rows

# --- Check if any row has real data (excluding ID / S.No) ---
def has_real_data(rows, columns):
    for row in rows:
        for col in columns:
            if col in ["ID", "S.No"]:
                continue
            val = row.get(col, {})
            if isinstance(val, dict):
                val = val.get("value", "")
            if isinstance(val, list):
                val = [v.get("value", v) if isinstance(v, dict) else v for v in val]
                val = [v for v in val if v not in (None, "")]
            if val not in (None, "", []):
                return True
    return False

from pptx.dml.color import RGBColor
from pptx.util import Pt

def update_overall_status_from_milestones(slide, key_milestones_data):
    """
    Updates the 'Overall Project Status' cell by appending the computed
    RAG status (GREEN, AMBER, RED) while forcing a white background and blue borders.
    """
    rag_values = []
    for row in key_milestones_data:
        rag = str(row.get("RAG", {}).get("value", "")).strip().upper()
        if rag in {"G", "R", "A"}:
            rag_values.append(rag)

    if not rag_values:
        return

    # Compute overall status (using pastel colors)
    if all(v == "G" for v in rag_values):
        overall_status, color, text_color = "GREEN", RGBColor(198, 239, 206), RGBColor(0, 97, 0)
    elif any(v == "R" for v in rag_values):
        overall_status, color, text_color = "RED", RGBColor(255, 199, 206), RGBColor(156, 0, 6)
    else:
        overall_status, color, text_color = "AMBER", RGBColor(255, 235, 156), RGBColor(156, 101, 0)

    # Find Table 3
    tbl_shape = None
    for shape in slide.shapes:
        if shape.name.lower().strip() == "table 3" and shape.has_table:
            tbl_shape = shape
            break

    if not tbl_shape:
        return

    table = tbl_shape.table
    clear_table_style(table)

    # Style Table 3: White background + Prolifics Blue border
    for row in table.rows:
        for cell in row.cells:
            set_cell_background(cell, "FFFFFF")
            set_cell_border(cell, border_color="00AEEF", border_width="12700") 
            cell.margin_left = int(Inches(0.1))
            cell.margin_right = int(Inches(0.1))
            cell.margin_top = int(Inches(0.1))

    # Set Slide 2 layout for Table 3
    tbl_shape.left = int(Inches(0.5))
    tbl_shape.top = int(Inches(1.4))
    tbl_shape.width = int(Inches(5.0))
    if len(table.columns) > 0:
        table.columns[0].width = int(Inches(4.8))

    # Locate the cell with 'Overall Project Status'
    for row in table.rows:
        for cell in row.cells:
            if "overall project status" in cell.text.lower():
                tf = cell.text_frame
                # Clear runs and set custom status text
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    p.text = ""
                    
                    r1 = p.add_run()
                    r1.text = "Overall Project Status is "
                    r1.font.bold = True
                    r1.font.size = Pt(14)
                    r1.font.color.rgb = RGBColor(0, 0, 0)
                    
                    r2 = p.add_run()
                    r2.text = overall_status
                    r2.font.bold = True
                    r2.font.size = Pt(14)
                    r2.font.name = "Calibri"
                    
                    # Highlight the status text color
                    r2.font.color.rgb = text_color if text_color else color
                    p.alignment = PP_ALIGN.LEFT
                    
                    # Ensure the cell background is WHITE as per screenshot
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

                # Format labels (Why?, Summary, etc.) in Prolifics Blue
                for para in tf.paragraphs:
                    txt = para.text.strip().lower()
                    para.alignment = PP_ALIGN.LEFT
                    for run in para.runs:
                        run.font.name = "Calibri"
                        run.font.size = Pt(12)
                        if any(kw in txt for kw in ["why", "summary", "planned", "financial"]):
                            run.font.color.rgb = RGBColor(0, 174, 239)
                            run.font.underline = True
                            run.font.bold = True
                        elif "overall project status" not in txt:
                            run.font.color.rgb = RGBColor(0, 0, 0)
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                return

    print("⚠ Could not find 'Overall Project Status' text in Table 3.")


def fill_roles_in_table(slide, project_roles, table_name="Table 5"):
    """
    Finds a table by name (e.g., 'Table 5') and fills in the corresponding
    names from the project_roles dictionary.
    Each role in the first column gets its name filled in the second column.
    """
    if not project_roles:
        return

    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    # Find the table by its name
    target_table = None
    for shape in slide.shapes:
        if shape.name.strip().lower() == table_name.strip().lower() and shape.has_table:
            target_table = shape.table
            break

    if not target_table:
        print(f"Table '{table_name}' not found on this slide.")
        return

    # Iterate over table rows (skip header if applicable)
    for row in target_table.rows:
        role_text = row.cells[0].text.strip()
        if role_text in project_roles and project_roles[role_text]:
            name_text = project_roles[role_text].strip()
            row.cells[1].text = name_text

            # Optional: style the text
            for paragraph in row.cells[1].text_frame.paragraphs:
                paragraph.font.name = "Calibri"
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)


from pptx.util import Pt
from copy import deepcopy

def safe_add_row(table):
    """
    Add a new row to a pptx table while preserving formatting.
    Returns the new row object (table.rows[-1]).
    """
    try:
        # For older python-pptx versions with no add_row() support
        tbl = table._tbl  # underlying XML element
        last_row = tbl.tr_lst[-1]
        new_row = deepcopy(last_row)
        tbl.append(new_row)
        return table.rows[-1]
    except Exception as e:
        print(f"Failed to add row safely: {e}")
        return None


def fill_roles_table(prs, project_roles):
    """
    Fill 'Table 5' on the first slide with project roles.
    The first row is treated as header, so data starts from the 2nd row.
    Dynamically adds or removes rows as needed.
    """

    # Validate and clean input
    if not project_roles or not isinstance(project_roles, dict):
        print("No valid roles to fill.")
        return

    # Remove empty or invalid entries
    filtered_roles = {
        role.strip(): person.strip()
        for role, person in project_roles.items()
        if role and role.strip() and person and person.strip()
    }

    if not filtered_roles:
        print("No non-empty roles to fill.")
        return

    # Locate the first slide and 'Table 5'
    first_slide = prs.slides[0]
    table_shape = None

    for shape in first_slide.shapes:
        if shape.name.strip().lower() == "table 5":
            table_shape = shape
            break

    if not table_shape:
        print("Could not find shape named 'Table 5'.")
        return

    table = table_shape.table

    # Adjust rows to fit exactly the number of roles (plus 1 for header)
    existing_rows = len(table.rows)
    required_rows = len(filtered_roles) + 1  # header row stays

    # Add missing rows
    while existing_rows < required_rows:
        safe_add_row(table)
        existing_rows += 1

    # Remove extra rows
    while existing_rows > required_rows:
        tbl = table._tbl  # underlying XML
        tbl.remove(tbl.tr_lst[-1])  # remove last row
        existing_rows -= 1

    # Fill data starting from 2nd row (index 1)
    for i, (role, person) in enumerate(filtered_roles.items(), start=1):
        role_cell = table.cell(i, 0)
        person_cell = table.cell(i, 1)

        role_cell.text = person   # Column 0: Name
        person_cell.text = role     # Column 1: Designation/Role

        # Apply text formatting
        for cell in (role_cell, person_cell):
            for p in cell.text_frame.paragraphs:
                if p.runs:
                    for run in p.runs:
                        run.font.size = Pt(12)
                        run.font.name = "Calibri"
                else:
                    p.font.size = Pt(12)
                    p.font.name = "Calibri"

        print(f"Filled row {i + 1}: {role} → {person}")

    print(f" Table 5 filled successfully with {len(filtered_roles)} roles.")


def fill_sb_architects_table(prs, project_roles):
    """
    Finds 'SWB Head of Customer Solutions:' or 'SB Architects:' inside Table 15 
    on slide 2, renames it to 'SB Architects:', and inserts joined architect names.
    """
    try:
        if len(prs.slides) < 2:
            print("No 2nd slide found.")
            return

        slide = prs.slides[1]
        target_table = None

        # Find Table 15
        for shape in slide.shapes:
            if shape.name.strip().lower() == "table 15":
                target_table = shape.table
                break

        if not target_table:
            print("Table 15 not found on 2nd slide.")
            return

        # Determine SB Architects names
        architect_names = []
        # Priority: roles containing 'architect' but NOT 'prolifics'
        for role, name in project_roles.items():
            if not role or not name:
                continue
            r_low = role.lower()
            if "architect" in r_low and "prolifics" not in r_low:
                n = name.strip()
                if n and n not in architect_names:
                    architect_names.append(n)
        
        display_name = " / ".join(architect_names) if architect_names else "N/A"

        # Locate label in Table 15
        for row_idx, row in enumerate(target_table.rows):
            for col_index, cell in enumerate(row.cells):
                txt = cell.text.lower()
                if "swb head" in txt or "architects" in txt:
                    # 1. Rename Label to 'SB Architects:'
                    cell.text = "SB Architects:"
                    for p in cell.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.LEFT
                        for r in p.runs:
                            r.font.color.rgb = RGBColor(255, 255, 255)
                            r.font.bold = True
                            r.font.size = Pt(12)
                            r.font.name = "Calibri"

                    # 2. Fill Name in next cell
                    if col_index + 1 < len(row.cells):
                        next_cell = row.cells[col_index + 1]
                        next_cell.text = display_name
                        for p in next_cell.text_frame.paragraphs:
                            p.alignment = PP_ALIGN.CENTER
                            for r in p.runs:
                                r.font.color.rgb = RGBColor(255, 255, 255)
                                r.font.bold = True
                                r.font.size = Pt(12)
                                r.font.name = "Calibri"
                        
                        # Ensure both cells are blue
                        set_cell_background(cell, "00AEEF")
                        set_cell_background(next_cell, "00AEEF")
                        set_cell_border(cell, border_color="00AEEF", border_width="12700")
                        set_cell_border(next_cell, border_color="00AEEF", border_width="12700")

                        # Set Slide 2 layout for Table 15 (Architects) to match Table 4/18
                        shape.left = int(5.6 * 914400)
                        shape.width = int(7.2 * 914400)
                        
                        # Set column widths for Table 15 to match Table 18 layout
                        if len(target_table.columns) > 1:
                            target_table.columns[0].width = int(Inches(3.0))
                            target_table.columns[1].width = int(Inches(4.2))

                        # Try to find Table 18 to align with it
                        ref_top = int(1.4 * 914400) + int(Inches(0.4)) # fallback top
                        for s in slide.shapes:
                            if s.name.strip().lower() == "table 18":
                                ref_top = s.top + s.height + int(Inches(0.08)) # margin below Table 18
                                break
                        
                        shape.top = ref_top

                        print(f"Updated SB Architects: {display_name} and moved it below Table 18")
                        return
                    else:
                        print("No next cell found to insert Project Manager name.")
                        return

        print("'SWB Head of Customer Solutions:' not found in Table 15.")

    except Exception as e:
        print(f"Error inserting Project Manager name: {e}")


def merge_dependencies_table(table, start_row=1, skip_column="Comments"):
    """
    Merge cells downward in Dependencies table for rows where:
    - Only one column has value per row except skip_column
    - Do NOT merge the skip_column
    - Skip headers and already merged cells
    """
    num_rows = len(table.rows)
    num_cols = len(table.columns)

    # Find index of skip_column
    skip_idx = None
    for c in range(num_cols):
        if table.cell(0, c).text.strip() == skip_column:
            skip_idx = c
            break

    for r in range(start_row, num_rows):
        # Identify all non-blank cells except skip_column
        non_blank_cells = [c for c in range(num_cols) 
                           if c != skip_idx and table.cell(r, c).text.strip() != ""]
        if len(non_blank_cells) != 1:
            continue  # skip row if multiple non-blank cells

        col_idx = non_blank_cells[0]
        merge_start = r

        # Merge downward as long as cells below in the same column are blank
        for rr in range(r + 1, num_rows):
            below_cell = table.cell(rr, col_idx)
            if hasattr(below_cell, 'merged') and below_cell.merged:
                break
            if below_cell.text.strip() != "":
                break

        merge_end = rr - 1
        if merge_end > merge_start:
            try:
                table.cell(merge_start, col_idx).merge(table.cell(merge_end, col_idx))
            except ValueError:
                pass

def is_long_content_table(rows):
    LONG_CONTENT_THRESHOLD = 200  # characters
    for row in rows:
        for cell in row.values():
            value = cell.get("value") if isinstance(cell, dict) else str(cell)
            # Convert lists/dicts to string
            if isinstance(value, list):
                value = "\n".join(str(v.get("value", v)) if isinstance(v, dict) else str(v) for v in value)
            if len(str(value)) > LONG_CONTENT_THRESHOLD:
                return True
    return False


def estimate_row_height(text, font_size_pt=10, char_per_line=55):
    if not text:
        return int(Pt(9).pt * 12700) # Even thinner base

    lines = 0
    for line in str(text).split("\n"):
        lines += 1 + len(line) // char_per_line

    # Safer multiplier 1.25 to prevent overlap and accommodate wrapping
    height_pt = font_size_pt * 1.25 * lines 
    height_emu = int(height_pt * 12700)

    # Maximum safety height
    MAX_ROW_HEIGHT_EMU = int(Inches(1.5))
    return min(height_emu, MAX_ROW_HEIGHT_EMU)



def is_visible_row(row, columns):
    """
    Returns True if the row should be included in the PPT table.

    Rules:
    1. Always include section titles (row has 'section_title').
    2. Exclude rows explicitly hidden (row.get("hidden") == True).
    3. Include data rows with at least one non-empty value in the given columns.
    """
    # Always keep section title rows
    if "section_title" in row and row["section_title"] not in (None, ""):
        return True

    # Skip hidden rows
    if row.get("hidden") is True:
        return False

    # Keep row if at least one column has data
    for col in columns:
        value = row.get(col, {})
        if isinstance(value, dict):
            val = value.get("value")
        else:
            val = value
        if val not in (None, "", []):
            return True

    return False

def merge_first_column_duplicates(table):
    """
    Merge consecutive rows in the first column if they have the same value.
    Only merges consecutive duplicates.
    Works dynamically for whatever the first column is.
    """
    if not table.rows:
        return
    
    first_col_idx = 0  # always the first column
    prev_text = table.cell(0, first_col_idx).text.strip()
    start_merge = 0

    for r in range(1, len(table.rows)):
        current_text = table.cell(r, first_col_idx).text.strip()
        if current_text == prev_text and current_text != "":
            table.cell(r, first_col_idx).text = ""
        else:
            if r - 1 > start_merge:
                try:
                    table.cell(start_merge, first_col_idx).merge(table.cell(r - 1, first_col_idx))
                except ValueError:
                    pass
            start_merge = r
            prev_text = current_text

    # Merge last group if needed
    if len(table.rows) - 1 > start_merge:
        try:
            table.cell(start_merge, first_col_idx).merge(table.cell(len(table.rows) - 1, first_col_idx))
        except ValueError:
            pass