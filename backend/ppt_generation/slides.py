from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import pandas as pd

from utils.llm_helper import LLMHelper
from ppt_generation.ppt_helper import (
    LEFT_MARGIN, MAX_TABLE_WIDTH, SLIDE_HEIGHT,
    insert_slide, copy_shapes_from_slide,
    apply_table_borders, style_header_cell, style_data_cell,
    update_project_title, update_slide_date, add_project_status_title,
    calculate_content_based_widths, detect_date_columns, is_single_word, style_section_title_cell, get_project_start_date, fill_sb_architects_table,
    filter_valid_data_rows, update_project_start_date, update_key_milestones_table, update_overall_status_from_milestones, fill_roles_table, merge_dependencies_table,
    is_long_content_table, estimate_row_height, is_visible_row, merge_first_column_duplicates, update_footer_year,
    set_cell_background, set_cell_border
)
from ppt_generation.sheet_config import get_sheet_config


def _write_phase_subheader_row(table, row_cursor, columns, row_height):
    """
    Writes the 'From SOW | Utilised | Remaining' sub-header row for Phase tables.
    Classifies columns into groups, merges cells, and applies matching colors.
    """
    col_l = [c.strip().lower() for c in columns]

    sow_indices, util_indices, remain_indices, other_indices = [], [], [], []
    for i, c in enumerate(col_l):
        if "used" in c or "cost" in c:
            util_indices.append(i)
        elif "remaining" in c or "balance" in c:
            remain_indices.append(i)
        elif "days" in c or "rate" in c or "charges" in c:
            sow_indices.append(i)
        else:
            other_indices.append(i)

    # (label, start_idx, end_idx, bg_hex)
    spans = []
    if other_indices:
        # Cell above 'Services' should be white
        spans.append(("", other_indices[0], other_indices[-1], "FFFFFF"))
    if sow_indices:
        spans.append(("From SOW", sow_indices[0], sow_indices[-1], "DDEBF7"))
    if util_indices:
        spans.append(("Utilised", util_indices[0], util_indices[-1], "FFF2CC"))
    if remain_indices:
        spans.append(("Remaining", remain_indices[0], remain_indices[-1], "E2EFDA"))

    # Initialise all cells in this row as blank / white
    for col_idx in range(len(columns)):
        cell = table.cell(row_cursor, col_idx)
        cell.text = ""
        set_cell_background(cell, "DDEBF7")
        set_cell_border(cell)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    # Apply each span
    for (label, start_idx, end_idx, bg_hex) in spans:
        start_cell = table.cell(row_cursor, start_idx)
        if end_idx > start_idx:
            try:
                for ci in range(start_idx + 1, end_idx + 1):
                    table.cell(row_cursor, ci).text = ""
                start_cell.merge(table.cell(row_cursor, end_idx))
            except Exception:
                pass
        set_cell_background(start_cell, bg_hex)
        set_cell_border(start_cell)
        tf = start_cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if label:
            run = p.add_run()
            run.text = label
            run.font.bold = True
            run.font.size = Pt(8)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(0, 0, 0)
        start_cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    table.rows[row_cursor].height = row_height

def create_ppt_from_excel(template_path, excel_data, output_path, project_title="", project_roles=None, skip_header_slides=False, author="", summary_data=None):
    prs = Presentation(template_path)

    # Use summary_data for slides 1 & 2 if provided, otherwise fallback to excel_data
    sum_data = summary_data if summary_data is not None else excel_data

    if not skip_header_slides:
        # Update title, date and author all in one block (Rectangle 6)
        if project_title:
            update_project_title(prs.slides[0], project_title, author=author)
            if len(prs.slides) > 1:
                add_project_status_title(prs.slides[1], project_title)
        update_slide_date(prs.slides[0])
        
        if project_roles:
            print("\nFilling Table 5 with roles from frontend...")
            fill_roles_table(prs, project_roles)
        else:
            print("No roles received.")

        
        # Update project start date FIRST so Architects table can align to it
        project_start_date = get_project_start_date(sum_data)
        if len(prs.slides) > 1:
            update_project_start_date(prs.slides[1], project_start_date)

        # Add roles under project title and date
        if project_roles:
            print("Filling roles in PPT:", project_roles)
            fill_sb_architects_table(prs, project_roles)
        
        # Update copyright year in footers
        update_footer_year(prs)

        key_milestones_data = []
        for sheet_name, tables in sum_data.items():
            if "key milestone" in sheet_name.lower():
                for table in tables:
                    key_milestones_data.extend(table.get("rows", []))

        # Fill Table 4 with Key Milestones
        if key_milestones_data:
            update_key_milestones_table(prs.slides[1], key_milestones_data)
            # Update Overall Project Status in Table 3 based on RAGs
            update_overall_status_from_milestones(prs.slides[1], key_milestones_data)
    else:
        print("\n[APPEND MODE] Skipping header slides — appending data from additional Excel file")

    # Layout constants
    TOP_MARGIN_IN = 0.55
    BOTTOM_BUFFER_IN = 1.0
    FOOTER_HEIGHT_IN = 0.6
    RIGHT_MARGIN_IN = 0.25
    HEADER_ROW_HEIGHT_IN = 0.30  # FIX: tighter height for header rows
    TOP_MARGIN_EMU = int(Inches(TOP_MARGIN_IN))
    BOTTOM_BUFFER_EMU = int(Inches(BOTTOM_BUFFER_IN))
    RIGHT_MARGIN_EMU = int(Inches(RIGHT_MARGIN_IN))
    FOOTER_HEIGHT_EMU = int(Inches(FOOTER_HEIGHT_IN))
    SLIDE_WIDTH_EMU = int(prs.slide_width)
    SLIDE_HEIGHT_EMU = int(prs.slide_height)
    TABLE_TITLE_HEIGHT_EMU = int(Inches(0.32))
    HEADER_ROW_HEIGHT_EMU = int(Inches(0.28))
    TABLE_GAP_EMU = int(Inches(0.20))  # Bring tables closer

    usable_height_emu = SLIDE_HEIGHT_EMU - TOP_MARGIN_EMU - BOTTOM_BUFFER_EMU - FOOTER_HEIGHT_EMU
    # When appending (skip_header_slides=True), insert new slides at the end
    base_insert_index = len(prs.slides) if skip_header_slides else 2
    remaining_space_emu = usable_height_emu
    current_slide = None
    added_so_far = 0
    

    # Prioritize Key Milestones and Deliverables to come immediately after Slide 2
    priority_sheets = ["key milestones", "key milestone", "deliverables", "deliverable", "risk", "fs", "financial", "phase", "decision", "action item"]
    def get_sheet_priority(s):
        s_lower = s.lower()
        if "key milestone" in s_lower:
            return 0
        if "deliverable" in s_lower:
            return 1
        if "risk" in s_lower:
            return 2
        if "assumption" in s_lower:
            return 3
        if "issue" in s_lower:
            return 4
        if "dependenc" in s_lower:
            return 5
        if "decision" in s_lower:
            return 6
        if "action item" in s_lower:
            return 7
        return 10

    sorted_sheet_names = sorted(
        excel_data.keys(),
        key=lambda s: (get_sheet_priority(s), s.lower())
    )

    for sheet_name in sorted_sheet_names:
        tables = excel_data[sheet_name]
        
        # Get per-sheet config (FS, Phase get custom widths, etc.)
        sheet_cfg = get_sheet_config(sheet_name)
        
        COMPACT_ROW_HEIGHT_IN = sheet_cfg["compact_row_height_in"]
        COMPACT_ROW_HEIGHT_EMU = sheet_cfg["compact_row_height_emu"]

        if not tables:
            continue

        # ── These sheets must always have their own dedicated slide(s) ──
        # Force a new slide so they never share with other sheets' tables
        is_solo_sheet = any(p in sheet_name.lower() for p in priority_sheets)
        if is_solo_sheet:
            current_slide = None
            remaining_space_emu = usable_height_emu

        sheet_slide_count = 1  # counts slides for this sheet
        added_titles = set()
        print(sheet_name)

        for table_data in tables:
            columns = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            if not columns:
                continue

            # -----------------------------
            # Expand rows with multiple Description values
            expanded_rows = []
            for row in rows:
                list_columns = {}
                for col, col_val in row.items():
                    cell_value = col_val.get("value") if isinstance(col_val, dict) else col_val
                    if isinstance(cell_value, list) and len(cell_value) > 1:
                        list_columns[col] = cell_value

                if not list_columns:
                    new_row = {col: col_val.copy() if isinstance(col_val, dict) else {"value": col_val, "style": {}}
                               for col, col_val in row.items()}
                    expanded_rows.append(new_row)
                    continue

                max_len = max(len(vals) for vals in list_columns.values())
                for i in range(max_len):
                    new_row = {}
                    for col, col_val in row.items():
                        cell = col_val.copy() if isinstance(col_val, dict) else {"value": col_val, "style": {}}
                        if col in list_columns:
                            value_list = list_columns[col]
                            if i < len(value_list):
                                val = value_list[i]
                                if isinstance(val, dict) and "value" in val:
                                    cell["value"] = val["value"]
                                else:
                                    cell["value"] = str(val)
                            else:
                                cell["value"] = ""
                        else:
                            if i > 0:
                                cell["value"] = ""
                        if not isinstance(cell, dict):
                            cell = {"value": cell, "style": {}}
                        new_row[col] = cell
                    expanded_rows.append(new_row)

            rows = [r for r in expanded_rows if is_visible_row(r, columns)]
            start_row = 0
            slide_count = 1

            # -----------------------------
            # Handle case: if rows are empty, show "None" instead of table
            if not rows:
                # Create a new slide if not already present
                if not current_slide:
                    slide_layout = prs.slide_layouts[1]
                    current_slide = insert_slide(prs, slide_layout, base_insert_index + added_so_far)
                    added_so_far += 1
                    remaining_space_emu = usable_height_emu

                # Determine top position for table title
                if remaining_space_emu == usable_height_emu:  # first content on slide
                    title_top = TOP_MARGIN_EMU
                else:
                    title_top = usable_height_emu - remaining_space_emu + TOP_MARGIN_EMU

                # Add table title
                if sheet_name in added_titles:
                    table_title_text = f"{sheet_name} (Contd..)"
                else:
                    table_title_text = sheet_name
                    added_titles.add(sheet_name)

                table_title_box = current_slide.shapes.add_textbox(
                    Inches(0.5),
                    title_top,
                    Inches(3.20),
                    TABLE_TITLE_HEIGHT_EMU
                )
                tf_table_title = table_title_box.text_frame
                tf_table_title.clear()
                p_table_title = tf_table_title.paragraphs[0]
                p_table_title.alignment = PP_ALIGN.LEFT
                run = p_table_title.add_run()
                run.text = table_title_text
                run.font.name = "Tahoma"
                run.font.bold = True
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(20, 54, 107)

                # Add "None" below table title
                none_box = current_slide.shapes.add_textbox(
                    Inches(0.7),
                    title_top + TABLE_TITLE_HEIGHT_EMU + Inches(0.2),
                    Inches(4),
                    Inches(0.3)
                )
                tf_none = none_box.text_frame
                tf_none.clear()
                p_none = tf_none.paragraphs[0]
                run_none = p_none.add_run()
                run_none.text = u"• None"
                run_none.font.name = "Calibri"
                run_none.font.size = Pt(12)
                run_none.font.italic = True
                run_none.font.color.rgb = RGBColor(89, 89, 89)
                p_none.alignment = PP_ALIGN.LEFT

                # Add consistent bottom spacing for 'None' sections too
                remaining_space_emu -= (TABLE_TITLE_HEIGHT_EMU + int(Inches(0.5)) + int(Inches(0.3)))
                remaining_space_emu = max(0, remaining_space_emu)
                continue


            while start_row < len(rows):
                use_dynamic_height = is_long_content_table(rows)

                # FIX #1: Always use remaining_space_emu for available space calculation.
                # Previously, slides < 4 ignored remaining space which caused overlaps
                # when multiple tables were placed on the same slide.
                available_space_emu = remaining_space_emu - TABLE_TITLE_HEIGHT_EMU - int(Inches(0.2))

                # Also account for header row in available space
                available_space_emu -= HEADER_ROW_HEIGHT_EMU

                COMPACT_ROW_HEIGHT_EMU = int(Inches(0.16))
                rows_to_fit = []
                current_height = 0

                for r in rows[start_row:]:
                    # Estimate row height
                    if use_dynamic_height:
                        row_height = max(
                            estimate_row_height(str(cell.get("value", "")), font_size_pt=12)
                            for cell in r.values()
                        )
                    else:
                        row_height = COMPACT_ROW_HEIGHT_EMU

                    # Stop if this row would overflow
                    if current_height + row_height > available_space_emu:
                        break

                    rows_to_fit.append(r)
                    current_height += row_height

                # If no rows can fit, create a new slide
                if not rows_to_fit:
                    slide_layout = prs.slide_layouts[1]
                    current_slide = insert_slide(prs, slide_layout, base_insert_index + added_so_far)
                    added_so_far += 1
                    remaining_space_emu = usable_height_emu
                    continue

                # FIX #2: Include header row height in total table height for accurate tracking.
                # Previously only data row heights were counted, but the table also renders
                # a header row, so remaining space was under-reported.
                table_height_emu = current_height + HEADER_ROW_HEIGHT_EMU

                # If this is the first table or title won't fit -> create new slide
                if not current_slide or (TABLE_TITLE_HEIGHT_EMU + table_height_emu > remaining_space_emu):
                    slide_layout = prs.slide_layouts[1]
                    current_slide = insert_slide(prs, slide_layout, base_insert_index + added_so_far)
                    added_so_far += 1
                    remaining_space_emu = usable_height_emu
                    continue

                # Compute table positioning
                if remaining_space_emu == usable_height_emu:
                    title_top = TOP_MARGIN_EMU
                    table_top_emu = title_top + TABLE_TITLE_HEIGHT_EMU + int(Inches(0.2))
                else:
                    title_top = usable_height_emu - remaining_space_emu + TOP_MARGIN_EMU
                    table_top_emu = title_top + TABLE_TITLE_HEIGHT_EMU

                # --- Add table title ---
                title_exists_on_slide = any(
                    shape.has_text_frame and sheet_name.split(" (Contd..)")[0] in shape.text
                    for shape in current_slide.shapes
                )

                if not title_exists_on_slide:
                    display_title = sheet_name
                    if sheet_name.lower() == "fs":
                        display_title = "Financial Summary:"
                    elif sheet_name.lower() == "phase":
                        display_title = "Financial Summary (Detailed by Phase):"

                    if display_title in added_titles:
                        table_title_text = f"{display_title} (Contd..)"
                    else:
                        table_title_text = display_title
                        added_titles.add(display_title)

                    table_title_box = current_slide.shapes.add_textbox(
                        Inches(0.5),
                        title_top,
                        Inches(3.33),
                        TABLE_TITLE_HEIGHT_EMU
                    )
                    tf_table_title = table_title_box.text_frame
                    tf_table_title.clear()
                    p_table_title = tf_table_title.paragraphs[0]
                    p_table_title.alignment = PP_ALIGN.LEFT
                    run = p_table_title.add_run()
                    run.text = table_title_text
                    run.font.name = "Tahoma"
                    run.font.bold = True
                    run.font.size = Pt(20)
                    run.font.color.rgb = RGBColor(20, 54, 107)
                else:
                    print("Title already exists on slide, skipping addition.")
                    table_top_emu = title_top

                # -----------------------------
                # Group by section and add table rows
                section_groups, current_section, current_rows = [], None, []
                for r in rows_to_fit:
                    section_title = r.get("section_title")
                    if section_title not in (None, ""):
                        if current_section is not None or current_rows:
                            section_groups.append({"section_title": current_section, "rows": current_rows})
                        current_section = section_title
                        current_rows = []
                    else:
                        current_rows.append(r)
                if current_section is not None or current_rows:
                    section_groups.append({"section_title": current_section, "rows": current_rows})

                _is_phase = sheet_name.lower() == "phase"
                total_rows_needed = 0
                for g in section_groups:
                    if g["section_title"]:
                        total_rows_needed += 1
                    if g["rows"]:
                        valid_g = filter_valid_data_rows(g["rows"], columns)
                        if valid_g:
                            if _is_phase:
                                # Phase: sub-header row + col-header row per group
                                total_rows_needed += 2
                            elif start_row == 0 or section_groups.index(g) == 0:
                                # Other sheets: col-header only for first group
                                total_rows_needed += 1
                            total_rows_needed += len(valid_g)

                left_margin_emu = sheet_cfg["left_margin_emu"]
                right_margin_emu = sheet_cfg["right_margin_emu"]

                # Use per-sheet table width if configured, otherwise auto-calculate
                if sheet_cfg["table_width_emu"]:
                    table_width_emu = sheet_cfg["table_width_emu"]
                else:
                    table_width_emu = SLIDE_WIDTH_EMU - left_margin_emu - right_margin_emu

                table = current_slide.shapes.add_table(total_rows_needed, len(columns),
                                                      left_margin_emu, table_top_emu,
                                                      table_width_emu, table_height_emu).table

                # Disable table style so PowerPoint respects cell-level borders
                # table.style = None doesn't work — must clear XML directly
                from pptx.oxml.ns import qn as _qn
                _tbl_el = table._tbl
                _tblPr = _tbl_el.find(_qn('a:tblPr'))
                if _tblPr is not None:
                    _NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    _sid = _tblPr.find(_qn('a:tableStyleId'))
                    if _sid is not None:
                        _sid.text = '{00000000-0000-0000-0000-000000000000}'
                    else:
                        from pptx.oxml import parse_xml as _px
                        _sid_el = _px(f'<a:tableStyleId xmlns:a="{_NS}">' + '{00000000-0000-0000-0000-000000000000}</a:tableStyleId>')
                        _tblPr.append(_sid_el)
                    for _flag in ['firstRow','firstCol','lastRow','lastCol','bandRow','bandCol']:
                        _tblPr.set(_flag, '0')
                table.first_row = False
                table.horz_banding = False
                table.vert_banding = False


                for row in table.rows:
                    for cell in row.cells:
                        cell.margin_left = Inches(0.05)
                        cell.margin_right = Inches(0.05)
                        cell.margin_top = Inches(0.0)
                        cell.margin_bottom = Inches(0.0)
                        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

                col_widths = calculate_content_based_widths(columns, rows_to_fit, table_width_emu,
                                                            fixed_widths_in=sheet_cfg["fixed_col_widths_emu"])
                total_col_width = sum(col_widths.values())
                scale_factor = table_width_emu / total_col_width if total_col_width else 1.0
                for i, col_name in enumerate(columns):
                    table.columns[i].width = int(col_widths[col_name] * scale_factor)

                # ── HARD SAFETY CLAMP ──
                # Regardless of what happened above, NEVER let the table overflow the slide.
                # This is the last line of defense.
                actual_total = sum(c.width for c in table.columns)
                if actual_total > table_width_emu:
                    clamp_scale = table_width_emu / actual_total
                    for col in table.columns:
                        col.width = int(col.width * clamp_scale)
                    print(f"  [WARNING] SAFETY CLAMP: {sheet_name} table was {actual_total/914400:.1f}\" -> clamped to {sum(c.width for c in table.columns)/914400:.1f}\"")

                row_cursor = 0
                for group_idx, group in enumerate(section_groups):
                    # Section title
                    if group["section_title"]:
                        first_cell = table.cell(row_cursor, 0)
                        if len(columns) > 1:
                            last_cell = table.cell(row_cursor, len(columns) - 1)
                            first_cell.merge(last_cell)

                        section_val = group.get("section_title", "")
                        section_text = section_val.get("value", "") if isinstance(section_val, dict) else str(section_val)
                        section_style = section_val.get("style", {}) if isinstance(section_val, dict) else {}
                        style_section_title_cell(first_cell, section_text, section_style, sheet_name=sheet_name)
                        table.rows[row_cursor].height = COMPACT_ROW_HEIGHT_EMU if not use_dynamic_height else int(Inches(0.2))
                        row_cursor += 1

                    # Header row(s): Phase gets sub-header+col-header per group;
                    # other sheets get col-header only for the first group.
                    if group["rows"]:
                        write_header_now = _is_phase or (row_cursor <= 1)
                        if write_header_now:
                            if _is_phase:
                                # Sub-header: "From SOW | Utilised | Remaining"
                                _write_phase_subheader_row(table, row_cursor, columns, COMPACT_ROW_HEIGHT_EMU)
                                row_cursor += 1

                            # Column header row
                            for col_idx, col_name in enumerate(columns):
                                cell = table.cell(row_cursor, col_idx)
                                style_header_cell(cell, col_name, font_size=10, col_width_emu=table.columns[col_idx].width, sheet_name=sheet_name)
                            table.rows[row_cursor].height = HEADER_ROW_HEIGHT_EMU
                            row_cursor += 1

                        for row_data in filter_valid_data_rows(group["rows"], columns):
                            table.rows[row_cursor].height = COMPACT_ROW_HEIGHT_EMU if not use_dynamic_height else max(
                                estimate_row_height(str(cell.get("value", "")), font_size_pt=12)
                                for cell in row_data.values()
                            )

                            # Detect if this is a "Total" row
                            first_col_val = str(row_data.get(columns[0], {}).get("value", "")).lower()
                            is_total_row = "total" in first_col_val

                            for col_idx, col_name in enumerate(columns):
                                cell = table.cell(row_cursor, col_idx)
                                value = row_data.get(col_name, {})
                                if not isinstance(value, dict):
                                    value = {"value": value}

                                if isinstance(value.get("value"), list):
                                    new_vals = []
                                    for v in value["value"]:
                                        if isinstance(v, dict) and "value" in v:
                                            new_vals.append(str(v["value"]))
                                        elif v not in (None, ""):
                                            new_vals.append(str(v))
                                    value["value"] = "\n".join(new_vals)

                                style_data_cell(cell, value, sheet_name=sheet_name, header=col_name, is_total_row=is_total_row)
                            row_cursor += 1

                # Merge first-column duplicates only for non-financial sheets
                if sheet_name.lower() not in ["fs", "phase"]:
                    merge_first_column_duplicates(table)

                milestone_col_idx = columns.index("Milestone") if "Milestone" in columns else None
                if milestone_col_idx is not None:
                    start_merge = 0
                    prev_text = table.cell(0, milestone_col_idx).text.strip()
                    for r in range(1, row_cursor):
                        current_text = table.cell(r, milestone_col_idx).text.strip()
                        if current_text == prev_text and current_text != "":
                            table.cell(r, milestone_col_idx).text = ""
                        else:
                            if r - 1 > start_merge:
                                table.cell(start_merge, milestone_col_idx).merge(
                                    table.cell(r - 1, milestone_col_idx))
                            start_merge = r
                            prev_text = current_text
                    if row_cursor - 1 > start_merge:
                        table.cell(start_merge, milestone_col_idx).merge(
                            table.cell(row_cursor - 1, milestone_col_idx))
                
                # FS only: merge identical adjacent cells in the title row (row 0) only.
                # Phase handles its own sub-header merging via _write_phase_subheader_row.
                if sheet_name.lower() == "fs":
                    for r_idx in range(min(2, row_cursor)):
                        row = table.rows[r_idx]
                        c_idx = 0
                        while c_idx < len(table.columns) - 1:
                            txt = row.cells[c_idx].text.strip()
                            if not txt:
                                c_idx += 1
                                continue
                            span = 1
                            while c_idx + span < len(table.columns) and row.cells[c_idx + span].text.strip() == txt:
                                span += 1
                            if span > 1:
                                try:
                                    for s_idx in range(1, span):
                                        row.cells[c_idx + s_idx].text = ""
                                    row.cells[c_idx].merge(row.cells[c_idx + span - 1])
                                except Exception as e:
                                    print(f"Warning: FS merge failed at row {r_idx} col {c_idx}: {e}")
                                c_idx += span
                            else:
                                c_idx += 1

                # Special PO# merge for Summary table
                if sheet_name.lower() == "fs" and "PO#" in columns:
                    po_col_idx = columns.index("PO#")
                    data_start = 2 # title is 0, header is 1
                    total_row_idx = -1
                    for r in range(data_start, row_cursor):
                        if "total" in str(table.cell(r, 0).text).lower():
                            total_row_idx = r
                            break
                    
                    if total_row_idx > data_start + 1:
                        # Merge PO# from first data row to row before Total
                        try:
                            # CLEAR text of other cells in range to prevent repetition
                            for r in range(data_start + 1, total_row_idx):
                                table.cell(r, po_col_idx).text = ""
                                
                            table.cell(data_start, po_col_idx).merge(table.cell(total_row_idx - 1, po_col_idx))
                            # Ensure PO# text is centered and blue
                            cell = table.cell(data_start, po_col_idx)
                            for p in cell.text_frame.paragraphs:
                                p.alignment = PP_ALIGN.CENTER
                                for run in p.runs:
                                    run.font.bold = True
                                    run.font.color.rgb = RGBColor(0, 112, 192) # Professional Blue
                        except:
                            pass

                # Apply borders AFTER all population and merging is complete
                apply_table_borders(table, border_color="000000", border_width="12700")

                # Copy template shapes (RAG legend) — skip for sheets that don't need it
                skip_rag_legend = sheet_cfg.get("skip_rag_legend", False)
                if len(prs.slides) > 1 and not hasattr(current_slide, "template_shapes_copied") and not skip_rag_legend:
                    copy_shapes_from_slide(prs.slides[1], current_slide, [
                        "Flowchart: Connector 13",
                        "Flowchart: Connector 3",
                        "Flowchart: Connector 5",
                        "TextBox 1"
                    ])
                    current_slide.template_shapes_copied = True

                # FIX #3: Always track remaining space consistently.
                # Previously, slides < 4 reset remaining_space to full usable_height,
                # so the second table on the same slide would overlap the first.
                # Now we always subtract the space consumed by this table + title + gap.
                consumed = table_height_emu + TABLE_TITLE_HEIGHT_EMU + TABLE_GAP_EMU
                remaining_space_emu -= consumed
                remaining_space_emu = max(0, remaining_space_emu)

                start_row += len(rows_to_fit)
                slide_count += 1
                print(added_titles)

        # ── After solo sheets, force the next sheet onto a fresh slide ──
        if is_solo_sheet:
            current_slide = None
            remaining_space_emu = usable_height_emu

    # ── Move "Thank You" and "Appendix" slides to the very end ──
    # These template slides may end up in the middle after data slides are inserted.
    # Detect them and reorder so they're always last (Thank You -> Appendix order).
    _move_trailing_slides_to_end(prs)

    # Update footer year from 2025 -> 2026 across all slides
    print("\nUpdating footer year 2025 -> 2026...")
    update_footer_year(prs, old_year="2025", new_year="2026")

    prs.save(output_path)
    return output_path


def _update_author(slide, author_name):
    """
    Update the Author line on slide 1. The template has a text box (Rectangle 6)
    with paragraphs: [Title, "Weekly Project Status Report", "Date: ...", "Author: ..."]
    The Author paragraph typically has 2 runs: "Author: " + "Name".
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            para_text = para.text.strip()
            if para_text.lower().startswith("author"):
                # Found the author paragraph — update the name part
                if len(para.runs) >= 2:
                    # Run 0 = "Author: ", Run 1 = name -> just update run 1
                    para.runs[1].text = author_name
                elif len(para.runs) == 1:
                    para.runs[0].text = f"Author: {author_name}"
                else:
                    para.text = f"Author: {author_name}"
                print(f"Updated author to: {author_name}")
                return
    print(f"Author paragraph not found on slide — skipping")


def _move_trailing_slides_to_end(prs):
    """
    Find slides whose content matches 'Thank You' or 'Appendix' patterns
    and move them to the end of the presentation in order:
      ... all data slides ... -> Thank You -> Appendix(es)
    """
    from lxml import etree

    slide_list = prs.slides._sldIdLst  # <p:sldIdLst> element

    thankyou_indices = []
    appendix_indices = []

    for i, slide in enumerate(prs.slides):
        # Skip first 2 slides (title + status) — they stay fixed
        if i < 2:
            continue

        slide_text = ""
        has_table = False
        is_image_only = True
        for shape in slide.shapes:
            if shape.has_table:
                has_table = True
                is_image_only = False
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip().lower()
                slide_text += " " + txt
                if txt:
                    is_image_only = False

        # Thank You slide: typically has only images (no text, no tables)
        # or contains "thank you" text
        if "thank" in slide_text or (is_image_only and not has_table and i > 1):
            thankyou_indices.append(i)
        elif "appendix" in slide_text:
            appendix_indices.append(i)

    # Collect sldId elements to move (thank you first, then appendix)
    indices_to_move = thankyou_indices + appendix_indices
    if not indices_to_move:
        return

    # Get the sldId elements in order
    sld_id_elements = list(slide_list)
    elements_to_move = [sld_id_elements[i] for i in indices_to_move]

    # Remove them from current position and append at end
    for elem in elements_to_move:
        slide_list.remove(elem)
        slide_list.append(elem)

    print(f"Moved slides to end: ThankYou={thankyou_indices}, Appendix={appendix_indices}")