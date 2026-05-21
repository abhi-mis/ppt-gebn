"""
fix_ppt.py — Standalone post-generation PPT fixer.

Usage:
    python fix_ppt.py report.pptx
    python fix_ppt.py path/to/any_generated.pptx

Output:
    Saves fixed file as <filename>-fixed.pptx in the same directory.

What it fixes:
    1. Tables whose actual row-height sum exceeds the shape height → resizes shape
    2. Shapes (titles, tables, text boxes) that overflow below slide bottom → shifts up or splits
    3. Shapes that overlap each other vertically on the same slide → re-stacks with proper gaps
    4. RAG legend text colliding with table headers → nudges table down or legend up
    5. "Decisions" / "• None" text boxes pushed past slide bottom → moves to previous open space or new slide
"""

import sys
import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ──────────────────────────────────────────────
# CONSTANTS
# ──────────────────────────────────────────────
SLIDE_BOTTOM_MARGIN_IN = 0.55       # keep shapes above this from bottom
MIN_GAP_IN = 0.15                   # minimum vertical gap between shapes
TABLE_TITLE_GAP_IN = 0.05           # gap between section title and its table
FOOTER_RESERVE_IN = 0.50            # reserve for "© 2025 Prolifics" footer area

# RAG legend shapes (copied from template on every content slide)
RAG_SHAPE_NAMES = {
    "Flowchart: Connector 13",  # R
    "Flowchart: Connector 3",   # A
    "Flowchart: Connector 5",   # G
}
RAG_TEXTBOX_SNIPPET = "Project at risk"  # identifies the RAG legend text box


# ──────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────
def emu(inches):
    return int(Inches(inches))


def to_in(emu_val):
    return emu_val / 914400


def shape_bottom(shape):
    """Actual bottom of shape — for tables uses sum of row heights."""
    if shape.has_table:
        t = shape.table
        h = sum(t.rows[r].height for r in range(len(t.rows)))
        return shape.top + h
    return shape.top + shape.height


def is_rag_legend(shape):
    if shape.name in RAG_SHAPE_NAMES:
        return True
    if shape.has_text_frame and RAG_TEXTBOX_SNIPPET in shape.text_frame.text:
        return True
    return False


def is_table_title(shape):
    """Section title text boxes like 'Key Milestones', 'Risks', 'Deliverables (Contd..)'"""
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    if not text:
        return False
    # RAG legend is NOT a table title
    if is_rag_legend(shape):
        return False
    # Slide 1 specific shapes
    if shape.name in ("TextBox 3", "Rectangle 6"):
        return False
    # Typically short bold titles — the generated ones are TextBox 1 / TextBox 15 / TextBox 16
    # and contain section names without "Red:" / "Amber:" / "Green:"
    if "Project at risk" in text or "recover/mitigate" in text:
        return False
    if "Project on track" in text:
        return False
    # Footer text
    if text.startswith("*See Milestones"):
        return False
    if "©" in text or "Prolifics" in text.split()[-1:]:
        return False
    return True


def is_none_text(shape):
    """Detects '• None' placeholder text boxes."""
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    return text in ("• None", "•None", "None", "- None")


def is_content_shape(shape):
    """Returns True for shapes that are part of the main content flow (tables, titles, none-text)."""
    if shape.has_table:
        return True
    if is_table_title(shape):
        return True
    if is_none_text(shape):
        return True
    return False


def get_content_shapes(slide):
    """Get all content shapes sorted by vertical position (top)."""
    shapes = []
    for shape in slide.shapes:
        if is_content_shape(shape):
            shapes.append(shape)
    shapes.sort(key=lambda s: s.top)
    return shapes


def get_rag_shapes(slide):
    """Get RAG legend shapes."""
    return [s for s in slide.shapes if is_rag_legend(s)]


def find_title_for_table(content_shapes, table_shape):
    """Find the title textbox immediately above a table."""
    table_idx = None
    for i, s in enumerate(content_shapes):
        if s is table_shape:
            table_idx = i
            break
    if table_idx is None or table_idx == 0:
        return None
    prev = content_shapes[table_idx - 1]
    if is_table_title(prev) and not prev.has_table:
        return prev
    return None


# ──────────────────────────────────────────────
# FIX 1: Correct table row + shape heights
# ──────────────────────────────────────────────
def _estimate_text_lines(text, col_width_emu, font_size_pt=12):
    """
    Estimate how many lines text will wrap to in a cell of given width.
    Uses approximate character width for Calibri at given size.
    """
    if not text:
        return 1
    # Approximate chars per inch for Calibri at 12pt ≈ 10-11 chars/inch
    # Account for cell margins (~0.1" total)
    usable_width_in = max(0.5, to_in(col_width_emu) - 0.1)
    chars_per_line = max(5, int(usable_width_in * 10))

    total_lines = 0
    for paragraph in text.split("\n"):
        if not paragraph.strip():
            total_lines += 1
            continue
        lines_needed = max(1, -(-len(paragraph) // chars_per_line))  # ceil division
        total_lines += lines_needed

    return max(1, total_lines)


def _calc_min_row_height(table, row_idx, font_size_pt=12):
    """
    Calculate the minimum row height needed to fit the tallest cell's text.
    """
    max_lines = 1
    for ci in range(len(table.columns)):
        text = table.cell(row_idx, ci).text or ""
        col_w = table.columns[ci].width
        lines = _estimate_text_lines(text, col_w, font_size_pt)
        if lines > max_lines:
            max_lines = lines

    # line height ≈ font_size * 1.3 (with some padding)
    line_height_emu = int(font_size_pt * 1.3 * 12700)
    # Add cell padding top + bottom (~0.04" each)
    padding_emu = emu(0.08)
    return max_lines * line_height_emu + padding_emu


def fix_table_row_heights(slide, slide_height_emu):
    """
    Expand table row heights to better fit their text content.
    Distributes available slide space fairly among all tables on the slide.
    Caps expansion so tables don't grow beyond available space.
    """
    # Collect all table shapes sorted by position
    table_shapes = sorted(
        [s for s in slide.shapes if s.has_table],
        key=lambda s: s.top
    )
    if not table_shapes:
        return

    # Count non-table content shapes (titles, text) to estimate their space
    non_table_height = 0
    for shape in slide.shapes:
        if shape.has_table:
            continue
        if is_content_shape(shape):
            non_table_height += shape.height
    
    # RAG legend takes space too
    rag_shapes = get_rag_shapes(slide)
    rag_bottom = max((s.top + s.height for s in rag_shapes), default=0)
    
    # Total available for all tables = slide height - RAG - titles - gaps - footer
    content_shapes = get_content_shapes(slide)
    num_gaps = max(1, len(content_shapes) - 1)
    gap_space = emu(0.20) * num_gaps  # estimate gaps
    top_reserve = max(emu(0.55), rag_bottom + emu(0.12)) if rag_bottom > 0 else emu(0.55)
    footer_reserve = emu(0.30)
    
    total_available_for_tables = slide_height_emu - top_reserve - non_table_height - gap_space - footer_reserve
    total_available_for_tables = max(emu(2.0), total_available_for_tables)  # minimum 2"
    
    # Calculate ideal heights for each table
    table_ideals = []
    table_currents = []
    for shape in table_shapes:
        t = shape.table
        nrows = len(t.rows)
        ideal_total = 0
        current_total = 0
        for ri in range(nrows):
            current_h = t.rows[ri].height
            min_h = _calc_min_row_height(t, ri)
            ideal_total += max(current_h, min_h)
            current_total += current_h
        table_ideals.append(ideal_total)
        table_currents.append(current_total)
    
    total_ideal = sum(table_ideals)
    total_current = sum(table_currents)
    
    # Determine scale factor
    if total_ideal <= total_available_for_tables:
        # Everything fits at ideal — expand fully
        scale = 1.0
    else:
        # Need to cap — how much extra can we give?
        extra_available = max(0, total_available_for_tables - total_current)
        extra_needed = max(1, total_ideal - total_current)
        scale = min(1.0, extra_available / extra_needed)
    
    # Apply to each table
    for ti, shape in enumerate(table_shapes):
        t = shape.table
        nrows = len(t.rows)
        for ri in range(nrows):
            current_h = t.rows[ri].height
            ideal_h = max(current_h, _calc_min_row_height(t, ri))
            extra = ideal_h - current_h
            if extra > 0:
                t.rows[ri].height = current_h + int(extra * scale)


def fix_table_shape_heights(slide):
    """Ensure table shape height matches actual sum of row heights."""
    for shape in slide.shapes:
        if shape.has_table:
            t = shape.table
            actual_h = sum(t.rows[r].height for r in range(len(t.rows)))
            if abs(shape.height - actual_h) > emu(0.05):
                shape.height = actual_h


def shrink_overflowing_cell_fonts(slide):
    """
    For each table cell where text would overflow the row height,
    reduce the font size until the text fits.
    This prevents PowerPoint from auto-expanding rows beyond set heights.
    
    Minimum font size: 7pt.
    """
    MIN_FONT_PT = 7
    
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        t = shape.table
        for ri in range(len(t.rows)):
            row_h = t.rows[ri].height
            
            for ci in range(len(t.columns)):
                cell = t.cell(ri, ci)
                text = cell.text or ""
                if not text.strip():
                    continue
                    
                col_w = t.columns[ci].width
                
                # Check if text fits at current font size
                current_font_pt = 12  # default
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.size:
                            current_font_pt = run.font.size.pt
                            break
                    break
                
                # Calculate min height needed at current font
                lines = _estimate_text_lines(text, col_w, current_font_pt)
                line_h = int(current_font_pt * 1.3 * 12700)
                needed_h = lines * line_h + emu(0.08)
                
                if needed_h <= row_h:
                    continue  # fits fine
                
                # Shrink font until it fits
                new_size = current_font_pt
                while new_size > MIN_FONT_PT:
                    new_size -= 0.5
                    # Recalculate with smaller font
                    # More chars per line at smaller size
                    usable_w = max(0.5, to_in(col_w) - 0.1)
                    # chars/inch scales roughly linearly with 1/font_size
                    cpi = 10 * (12 / new_size)
                    cpl = max(5, int(usable_w * cpi))
                    
                    new_lines = 0
                    for para in text.split("\n"):
                        if not para.strip():
                            new_lines += 1
                            continue
                        new_lines += max(1, -(-len(para) // cpl))
                    
                    new_line_h = int(new_size * 1.3 * 12700)
                    new_needed = new_lines * new_line_h + emu(0.08)
                    
                    if new_needed <= row_h:
                        break
                
                if new_size < current_font_pt:
                    # Apply the reduced font size
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(new_size)


# ──────────────────────────────────────────────
# FIX 2: Re-stack content shapes to eliminate overlaps
# ──────────────────────────────────────────────
def restack_content_shapes(slide, slide_height_emu):
    """
    Re-position content shapes (titles + tables) so they flow top-to-bottom
    with proper gaps and no overlaps.
    Uses adaptive footer reserve — starts generous, shrinks if needed.
    """
    content = get_content_shapes(slide)
    if not content:
        return

    rag_shapes = get_rag_shapes(slide)
    rag_bottom = 0
    for rs in rag_shapes:
        b = rs.top + rs.height
        if b > rag_bottom:
            rag_bottom = b

    # Content must start below RAG legend (if present)
    min_content_top = emu(0.55)  # default top margin
    if rag_bottom > 0:
        min_content_top = max(min_content_top, rag_bottom + emu(0.12))

    # Try multiple footer reserves — generous first, then tighter
    footer_reserves = [0.50, 0.30, 0.15, 0.05]
    
    for footer_in in footer_reserves:
        max_bottom = slide_height_emu - emu(footer_in)
        success = _try_layout(content, min_content_top, max_bottom)
        
        if success:
            return
    
    # Last resort: allow content up to absolute slide edge
    _try_layout(content, min_content_top, slide_height_emu)


def _try_layout(content, min_content_top, max_bottom):
    """
    Attempt to lay out content shapes within the given bounds.
    Returns True if everything fits, False if overflow.
    """
    cursor = min_content_top

    for i, shape in enumerate(content):
        if i > 0:
            prev = content[i - 1]
            if is_table_title(prev) and shape.has_table:
                cursor += emu(TABLE_TITLE_GAP_IN)
            elif prev.has_table and is_table_title(shape):
                cursor += emu(0.25)
            elif is_table_title(prev) and is_none_text(shape):
                # Title followed by "• None" — tight
                cursor += emu(0.08)
            elif is_none_text(prev):
                # After "• None", next section title
                cursor += emu(0.25)
            else:
                cursor += emu(MIN_GAP_IN)

        shape.top = cursor
        cursor += shape.height

    last_bottom = shape_bottom(content[-1]) if content else 0
    
    if last_bottom <= max_bottom:
        return True
    
    # Try compressing gaps
    overflow = last_bottom - max_bottom
    _try_compress(content, overflow, min_content_top, max_bottom)
    
    last_bottom = shape_bottom(content[-1]) if content else 0
    return last_bottom <= max_bottom


def _try_compress(content, overflow, min_top, max_bottom):
    """Try to reduce gaps to fit everything on the slide."""
    if len(content) <= 1:
        return

    # Calculate total gap space available to compress
    gaps = []
    for i in range(1, len(content)):
        gap = content[i].top - shape_bottom(content[i - 1])
        gaps.append(gap)

    total_gap = sum(gaps)
    min_total_gap = emu(0.02) * len(gaps)  # absolute minimum 0.02" per gap

    compressible = total_gap - min_total_gap
    if compressible <= 0:
        return  # can't compress further

    # Compress proportionally
    compress_ratio = min(1.0, overflow / compressible)

    cursor = content[0].top
    for i, shape in enumerate(content):
        shape.top = cursor
        if i < len(gaps):
            old_gap = gaps[i]
            min_gap = emu(0.02)
            reduction = int((old_gap - min_gap) * compress_ratio)
            new_gap = max(min_gap, old_gap - reduction)
            cursor = shape_bottom(shape) + new_gap
        else:
            cursor = shape_bottom(shape)


# ──────────────────────────────────────────────
# FIX 3: Handle shapes that overflow past slide bottom
# ──────────────────────────────────────────────
def fix_overflow_shapes(slide, slide_height_emu):
    """
    Move shapes that are completely or mostly below the slide boundary.
    Only handles non-content shapes (footers, stray elements).
    Content shapes are already handled by restack_content_shapes.
    """
    for shape in slide.shapes:
        if is_rag_legend(shape):
            continue
        # Skip content shapes — already positioned by restack
        if is_content_shape(shape):
            continue
            
        bot = shape.top + shape.height
        if shape.top >= slide_height_emu:
            # Completely off-slide — find the lowest content shape and put below it
            lowest = _find_lowest_content_bottom(slide)
            shape.top = lowest + emu(0.15)
        elif bot > slide_height_emu and not shape.has_table:
            # Non-content text box overflowing — nudge up
            overshoot = bot - slide_height_emu
            shape.top = max(emu(0.5), shape.top - overshoot)


def _find_lowest_content_bottom(slide):
    """Find the bottom of the lowest content shape on the slide."""
    lowest = 0
    for shape in slide.shapes:
        if is_rag_legend(shape):
            continue
        b = shape_bottom(shape)
        if b > lowest:
            lowest = b
    return lowest


# ──────────────────────────────────────────────
# FIX 4: Fix RAG legend vs table header collision
# ──────────────────────────────────────────────
def fix_rag_legend_collision(slide):
    """
    If RAG legend text overlaps with the first content shape,
    shift the RAG legend up or the content down.
    """
    rag_shapes = get_rag_shapes(slide)
    if not rag_shapes:
        return

    rag_bottom = max(s.top + s.height for s in rag_shapes)
    content = get_content_shapes(slide)
    if not content:
        return

    first_content_top = content[0].top
    if rag_bottom > first_content_top:
        # Overlap detected — shift all content down
        shift = rag_bottom - first_content_top + emu(0.10)
        for shape in content:
            shape.top += shift


# ──────────────────────────────────────────────
# MAIN FIXER
# ──────────────────────────────────────────────
def fix_ppt(input_path):
    prs = Presentation(input_path)
    slide_h = prs.slide_height

    # Skip slide 1 (title slide), slide 2 (template status slide),
    # slide 10 (thank you), slide 11 (appendix) — only fix generated content slides
    # But we still check all slides for safety

    total_slides = len(prs.slides)

    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        print(f"Processing slide {slide_num}/{total_slides}...")

        # Skip title slide (slide 1) — has its own custom layout
        if slide_num == 1:
            print(f"  Skipping slide 1 (title slide)")
            continue

        # Skip template status slide (slide 2) — has pre-built layout
        if slide_num == 2:
            print(f"  Skipping slide 2 (template status slide)")
            continue

        # Skip thank you slide (typically second-to-last or has only a picture)
        shape_types = [s.shape_type for s in slide.shapes]
        if len(slide.shapes) == 1 and slide.shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"  Skipping slide {slide_num} (thank you / image-only)")
            continue

        # Skip appendix slide (has PLACEHOLDER shapes)
        has_placeholder = any(
            s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER for s in slide.shapes
        )
        if has_placeholder:
            print(f"  Skipping slide {slide_num} (appendix / placeholder slide)")
            continue

        # ── Apply fixes in order ──
        print(f"  Fixing slide {slide_num}...")

        # Step 1a: Expand row heights to fit text content (prevents PowerPoint auto-expand)
        fix_table_row_heights(slide, slide_h)

        # Step 1b: Correct table shape heights to match actual row heights
        fix_table_shape_heights(slide)

        # Step 1c: Shrink fonts in cells that would overflow their row height
        shrink_overflowing_cell_fonts(slide)

        # Step 2: Fix RAG legend collision with content
        fix_rag_legend_collision(slide)

        # Step 3: Re-stack all content shapes with proper gaps
        restack_content_shapes(slide, slide_h)

        # Step 4: Fix any remaining overflow past slide bottom
        fix_overflow_shapes(slide, slide_h)

        # Verify
        _verify_slide(slide, slide_num, slide_h)

    # Build output path
    # ── Final pass: sanitizer catches anything the above fixes missed ──
    from ppt_sanitizer import sanitize_presentation
    print("\nRunning final layout sanitizer...")
    sanitize_presentation(prs)

    base, ext = os.path.splitext(input_path)
    output_path = f"{base}-fixed{ext}"
    prs.save(output_path)
    print(f"\nSaved fixed PPT: {output_path}")
    return output_path


def _verify_slide(slide, slide_num, slide_h):
    """Print warnings for any remaining issues."""
    max_bot = slide_h - emu(FOOTER_RESERVE_IN)
    for shape in slide.shapes:
        if is_rag_legend(shape):
            continue
        bot = shape_bottom(shape)
        if bot > slide_h:
            print(f"  ⚠ WARNING slide {slide_num}: '{shape.name}' bottom at {to_in(bot):.2f}in exceeds slide height {to_in(slide_h):.2f}in")

    # Check overlaps between content shapes
    content = get_content_shapes(slide)
    for i in range(len(content) - 1):
        a = content[i]
        b = content[i + 1]
        a_bot = shape_bottom(a)
        b_top = b.top
        if a_bot > b_top + emu(0.02):  # allow 0.02in tolerance
            a_name = a.text_frame.text[:30] if a.has_text_frame else a.name
            b_name = b.text_frame.text[:30] if b.has_text_frame else b.name
            print(f"  ⚠ WARNING slide {slide_num}: overlap between '{a_name}' (bot={to_in(a_bot):.2f}) and '{b_name}' (top={to_in(b_top):.2f})")


# ──────────────────────────────────────────────
# CLI ENTRY POINT
# ──────────────────────────────────────────────
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python fix_ppt.py <input.pptx>")
        print("Output: <input>-fixed.pptx")
        sys.exit(1)

    input_file = sys.argv[1]
    if not os.path.exists(input_file):
        print(f"Error: File not found: {input_file}")
        sys.exit(1)

    fix_ppt(input_file)
