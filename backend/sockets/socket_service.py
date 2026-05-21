import os
from utils.ppt_slide_summary import extract_slides_text
from pptx import Presentation
import json

from sockets.handlers import determine_slide_from_text
from ai_helper.open_ai_helper import OpenAIHelper


def process_event_message(data):
    print("Message received:", data)

    BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    prompt_file = os.path.join(BASE_DIR, "prompts", "event_message_prompt.txt")

    with open(prompt_file, "r", encoding="utf-8") as f:
        prompt_template = f.read()

    text = data.get("message", "")
    ppt_name = data.get("ppt_name", None)
    options = data.get("options", {})

    print("Options received:", options)

    if not text:
        print("No text provided in the message.")
        return "No message provided."

    if not ppt_name:
        print("No ppt name provided in the message.")
        return "No PPT file specified."

    ppt_path = os.path.join("generated", ppt_name)
    if not os.path.exists(ppt_path):
        print("PPT file not found:", ppt_path)
        return "PPT file not found."

    # -----------------------------------------
    # Determine slide number
    # -----------------------------------------
    slide_number = data.get("slideNumber", None)
    if slide_number is None:
        slide_number = determine_slide_from_text(text)
        if slide_number:
            print(f"LLM determined user wants slide {slide_number}")
        else:
            print("No specific slide mentioned -> processing full PPT")

    # -----------------------------------------
    # Extract slide text
    # -----------------------------------------
    if slide_number:
        try:
            slide_number = int(slide_number)
        except ValueError:
            slide_number = None

        slide_data = extract_ppt_text(ppt_path, slide_number)
        if not slide_data or "error" in slide_data:
            return "Slide not found."
        slide_text = slide_data["text"]
    else:
        slides = extract_ppt_text(ppt_path)
        slide_text = "\n\n".join(
            [f"Slide {s['slide_number']}:\n{s['text']}" for s in slides]
        )

    # -----------------------------------------
    # Build AI behavior from options
    # -----------------------------------------
    behavior = build_ai_behavior(options)

    # -----------------------------------------
    # Build and send prompt
    # -----------------------------------------
    ai = OpenAIHelper()

    base_prompt = (
        prompt_template
        .replace("{{text}}", slide_text)
        .replace("{{question}}", text)
    )

    final_prompt = build_prompt(base_prompt, text, behavior)
    print(f"\nFinal prompt length: {len(final_prompt)} chars")

    response = ai.ask(
        prompt=final_prompt,
        system_prompt="You are an expert PowerPoint analyst and project management assistant.",
        max_tokens=800,
        temperature=0.4
    )

    # -----------------------------------------
    # Post-processing based on options
    # -----------------------------------------
    if behavior.get("explanations"):
        explanation_text = (
            "Based on your settings: "
            + " ".join(behavior["explanations"])
        )
        response = f"{explanation_text}\n\n{response}"

    if "ask_for_feedback" in behavior.get("post_actions", []):
        response += "\n\nWould you like me to refine or improve this further?"

    print("Final Response ->", response[:200])
    return response


def build_ai_behavior(options):
    behavior = {
        "mode": "default",
        "instructions": [],
        "post_actions": [],
        "explanations": [],
    }

    if options.get("summary_per_sheet"):
        behavior["instructions"].append(
            "Generate summaries grouped by sheet/section, not by individual slide."
        )
        behavior["explanations"].append("Summaries grouped by sheet.")

    if options.get("jargons"):
        behavior["instructions"].append(
            "Identify technical terms, acronyms, and jargon. Provide brief plain-English explanations for each."
        )
        behavior["explanations"].append("Technical terms explained.")

    if options.get("risks_slide"):
        behavior["instructions"].append(
            "Identify any risks, warnings, blockers, or sensitive content. Flag items marked as Red/Amber RAG status."
        )
        behavior["explanations"].append("Risk detection enabled.")

    if options.get("recommend_improvements"):
        behavior["instructions"].append(
            "For any identified risks or issues, suggest specific actionable improvements."
        )
        behavior["explanations"].append("Improvement suggestions included.")

    if options.get("answer_style"):
        behavior["instructions"].append(
            "Keep the response concise, professional, and to-the-point. Use bullet points for clarity."
        )
        behavior["explanations"].append("Concise mode active.")

    if options.get("validate_ai"):
        behavior["instructions"].append(
            "Double-check facts against the slide content. If uncertain, say so explicitly."
        )
        behavior["explanations"].append("AI validation active.")

    if options.get("insert_summary"):
        behavior["post_actions"].append("insert_into_ppt")

    if options.get("feedback_loop"):
        behavior["post_actions"].append("ask_for_feedback")

    if not behavior["instructions"]:
        behavior["instructions"].append("Provide a clear and accurate response.")

    return behavior


def build_prompt(prompt_template, user_message, behavior):
    instructions = "\n".join(f"- {i}" for i in behavior["instructions"])

    return f"""USER REQUEST:
{user_message}

AI BEHAVIOR RULES:
{instructions}

TASK:
{prompt_template}
"""


def extract_ppt_text(ppt_path, slide_number=None):
    prs = Presentation(ppt_path)

    def extract_text_from_shape(shape):
        text_content = []
        if hasattr(shape, "text") and shape.text.strip():
            text_content.append(shape.text.strip())
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(" | ".join(row_text))
        return text_content

    if slide_number is not None:
        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": "Invalid slide number"}
        slide = prs.slides[slide_number - 1]
        slide_text = []
        for shape in slide.shapes:
            slide_text.extend(extract_text_from_shape(shape))
        return {"slide_number": slide_number, "text": "\n".join(slide_text)}

    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            slide_text.extend(extract_text_from_shape(shape))
        slides.append({"slide_number": idx, "text": "\n".join(slide_text)})
    return slides
