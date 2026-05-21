"""
refine_comments.py — Scans PPT for long comments (>20 words),
refines using AI, replaces in PPT, emits per-comment logs.
"""

import os
import traceback
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from ai_helper.open_ai_helper import OpenAIHelper

MIN_WORD_COUNT = 20
COMMENT_COLUMN_KEYWORDS = ["comment", "mitigation", "remark", "observation", "note"]
GENERATED_FOLDER = "./generated"


def is_comment_column(header_text):
    header_lower = header_text.lower().strip()
    return any(keyword in header_lower for keyword in COMMENT_COLUMN_KEYWORDS)


def get_cell_text(cell):
    parts = []
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                parts.append(run.text)
    return " ".join(parts).strip()


def word_count(text):
    if not text:
        return 0
    return len(text.split())


def replace_cell_text(cell, new_text):
    text_frame = cell.text_frame
    font_name = "Calibri"
    font_size = Pt(12)
    font_bold = False
    font_italic = False
    font_color = RGBColor(0, 0, 0)

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.name:
                font_name = run.font.name
            if run.font.size:
                font_size = run.font.size
            if run.font.bold is not None:
                font_bold = run.font.bold
            if run.font.italic is not None:
                font_italic = run.font.italic
            try:
                if run.font.color and run.font.color.rgb:
                    font_color = run.font.color.rgb
            except:
                pass
            break
        break

    text_frame.clear()
    lines = new_text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = font_size
        run.font.bold = font_bold
        run.font.italic = font_italic
        run.font.color.rgb = font_color


def scan_ppt_for_long_comments(prs):
    found_comments = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_table:
                continue
            table = shape.table
            num_rows = len(table.rows)
            num_cols = len(table.columns)
            if num_rows < 2:
                continue
            comment_col_indices = []
            for col_idx in range(num_cols):
                header_text = get_cell_text(table.cell(0, col_idx))
                if is_comment_column(header_text):
                    comment_col_indices.append((col_idx, header_text))
            if not comment_col_indices:
                continue
            for row_idx in range(1, num_rows):
                for col_idx, col_name in comment_col_indices:
                    cell = table.cell(row_idx, col_idx)
                    text = get_cell_text(cell)
                    wc = word_count(text)
                    if wc > MIN_WORD_COUNT:
                        found_comments.append({
                            "slide_idx": slide_idx,
                            "shape_idx": shape_idx,
                            "row_idx": row_idx,
                            "col_idx": col_idx,
                            "col_name": col_name,
                            "original_text": text,
                            "word_count": wc,
                        })
    return found_comments


def refine_comments_in_ppt(data, emit_fn=None):
    ppt_name = data.get("ppt_name", None)
    if not ppt_name:
        _emit(emit_fn, "refine_error", {"error": "No ppt_name provided.", "status": "error"})
        return

    ppt_path = os.path.join(GENERATED_FOLDER, ppt_name)
    if not os.path.exists(ppt_path):
        _emit(emit_fn, "refine_error", {"error": f"PPT file not found: {ppt_name}", "ppt_name": ppt_name, "status": "error"})
        return

    try:
        # Step 1: Scan
        _emit(emit_fn, "refine_status", {"status": "scanning", "message": "Scanning PPT for long comments..."})
        prs = Presentation(ppt_path)
        found_comments = scan_ppt_for_long_comments(prs)
        total = len(found_comments)
        print(f"\nFound {total} comments with > {MIN_WORD_COUNT} words")

        if total == 0:
            _emit(emit_fn, "refine_complete", {
                "status": "done", "refined_count": 0, "total_found": 0,
                "ppt_name": ppt_name,
                "message": "No comments found that need refinement (all are 20 words or fewer).",
            })
            return

        # Step 2: Refine
        _emit(emit_fn, "refine_status", {"status": "refining", "message": f"Refining {total} comments using AI...", "total": total})

        original_texts = [c["original_text"] for c in found_comments]
        ai_helper = OpenAIHelper()
        batch_size = 10
        refined_texts = []

        for batch_start in range(0, total, batch_size):
            batch_end = min(batch_start + batch_size, total)
            batch = original_texts[batch_start:batch_end]
            batch_refined = ai_helper.refine_comments_bulk(batch, batch_size=len(batch))
            refined_texts.extend(batch_refined)

            _emit(emit_fn, "refine_progress", {
                "total": total, "processed": len(refined_texts),
                "batch_info": f"Batch {batch_start + 1}-{batch_end} of {total}",
                "ppt_name": ppt_name,
            })

        # Step 3: Replace + emit per-comment log
        _emit(emit_fn, "refine_status", {"status": "replacing", "message": "Writing refined comments back into PPT..."})

        replaced_count = 0
        for i, comment_info in enumerate(found_comments):
            refined = refined_texts[i] if i < len(refined_texts) else ""
            original = comment_info["original_text"]

            if not refined or refined.startswith("Error:"):
                _emit(emit_fn, "refine_log", {
                    "index": i + 1, "total": total,
                    "slide": comment_info["slide_idx"] + 1,
                    "column": comment_info["col_name"],
                    "status": "skipped",
                    "reason": "AI returned empty/error",
                    "original": original[:100],
                    "refined": "",
                })
                continue

            if refined.strip() == original.strip():
                _emit(emit_fn, "refine_log", {
                    "index": i + 1, "total": total,
                    "slide": comment_info["slide_idx"] + 1,
                    "column": comment_info["col_name"],
                    "status": "unchanged",
                    "original": original[:100],
                    "refined": refined[:100],
                })
                continue

            slide = prs.slides[comment_info["slide_idx"]]
            shape = slide.shapes[comment_info["shape_idx"]]
            cell = shape.table.cell(comment_info["row_idx"], comment_info["col_idx"])
            replace_cell_text(cell, refined)
            replaced_count += 1

            _emit(emit_fn, "refine_log", {
                "index": i + 1, "total": total,
                "slide": comment_info["slide_idx"] + 1,
                "column": comment_info["col_name"],
                "status": "refined",
                "original": original[:120],
                "refined": refined[:120],
            })

        # Step 4: Sanitize layout
        _emit(emit_fn, "refine_status", {"status": "sanitizing", "message": "Fixing layout..."})
        from ppt_sanitizer import sanitize_presentation
        sanitize_presentation(prs)

        # Step 5: Save
        prs.save(ppt_path)
        print(f"\nSaved: {ppt_path} ({replaced_count}/{total} replaced)")

        _emit(emit_fn, "refine_complete", {
            "status": "done", "refined_count": replaced_count, "total_found": total,
            "ppt_name": ppt_name,
            "message": f"Refined {replaced_count} of {total} comments.",
        })

    except Exception as e:
        print(f"Error: {e}")
        traceback.print_exc()
        _emit(emit_fn, "refine_error", {"error": str(e), "ppt_name": ppt_name, "status": "error"})


def _emit(emit_fn, event, data):
    if emit_fn:
        emit_fn(event, data)
    print(f"  [{event}] {data.get('status', data.get('message', ''))}")
